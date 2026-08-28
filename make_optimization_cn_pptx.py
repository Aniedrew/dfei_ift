"""生成中文 PPT: DFEI 重建优化方案介绍 (2026-08-16)
覆盖: NoneIso 指标 / 问题诊断 / B2 可微剪枝 / 源检测头(RC训练化) / Loss再平衡 / LCA clip / 链级LCA过滤 / 评估矩阵
输出到 meeting_20260816_optim/DFEI_optimization_20260816.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE = "/lzufs/home/guoqingxiang/dfei/scalable_mtl_hgnn"
FIG = f"{BASE}/meeting_20260811_figs"
OUT_DIR = f"{BASE}/meeting_20260816_optim"
os.makedirs(OUT_DIR, exist_ok=True)
OUT = f"{OUT_DIR}/DFEI_optimization_20260816.pptx"

BLUE = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def add_title_bar(slide, text, color=BLUE):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = color
    ln = slide.shapes.add_shape(1, Inches(0.6), Inches(0.98), Inches(12.1), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = color; ln.line.fill.background()

def add_bullets(slide, items, left, top, width, height, size=16, color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = ("\u2022  " if isinstance(it, str) else "\u2022  ") + it
        r.font.size = Pt(size); r.font.color.rgb = color
    return tb

def add_pic(slide, fname, left, top, width=None, height=None):
    path = os.path.join(FIG, fname)
    if not os.path.exists(path):
        return
    from PIL import Image
    with Image.open(path) as im:
        w, h = im.size
    ar = w / h
    if width is None and height is None:
        width = Inches(7.0)
    if width is not None and height is None:
        height = width / ar
    if height is not None and width is None:
        width = height * ar
    slide.shapes.add_picture(path, left, top, width=width, height=height)

def add_footer(slide, idx):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.35))
    tf = tb.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = f"郭清祥 \u00b7 UCAS \u00b7 DFEI \u00b7 2026-08-16      {idx}"
    r.font.size = Pt(11); r.font.color.rgb = GRAY

def new_slide(idx):
    s = prs.slides.add_slide(BLANK)
    add_footer(s, idx)
    return s

def kv_rows(slide, rows, top=1.7, left=0.9, w=11.5, h=0.9, row_h=0.9, k_w=2.2, size=17):
    for i, (k, v) in enumerate(rows):
        y = Inches(top + row_h * i)
        tb = slide.shapes.add_textbox(Inches(left), y, Inches(k_w), Inches(h))
        tf = tb.text_frame
        r = tf.paragraphs[0].add_run(); r.text = k
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = BLUE
        tb = slide.shapes.add_textbox(Inches(left + k_w), y, Inches(w - k_w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = v
        r.font.size = Pt(size); r.font.color.rgb = DARK

# ============ Slide 1: 封面 ============
s = new_slide(1)
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.4))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "DFEI 重建优化方案介绍"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = BLUE
tb = s.shapes.add_textbox(Inches(1.0), Inches(3.7), Inches(11.3), Inches(0.9))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "从问题诊断到训练侧改进：B2 可微剪枝 \u00b7 源检测头 \u00b7 Loss 再平衡 \u00b7 链级 LCA 判据"
r.font.size = Pt(22); r.font.color.rgb = DARK
tb = s.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.7))
tf = tb.text_frame
r = tf.paragraphs[0].add_run()
r.text = "郭清祥 \u00b7 中国科学院大学 \u00b7 DFEI 组会 \u00b7 2026-08-16"
r.font.size = Pt(18); r.font.color.rgb = GRAY

# ============ Slide 2: NoneIso 是什么 ============
s = new_slide(2)
add_title_bar(s, "评估指标：NoneIso 是什么")
add_bullets(s, [
    "PerfectReco：重建出的链与 truth 完全一致（我们关心的核心指标）",
    "AllParticles：所有信号粒子都被重建出来（可能含多余粒子）",
    "NoneIso（非孤立）：重建链中混入了不属于该链的无关（背景）粒子 —— 链被污染",
    "PartReco：只重建出链的一部分；NotFound：完全没找到",
    "NoneIso 是 Perfect 之外最大的损失来源：v31 基线 NoneIso 高达 56.6%",
    "根因：剪枝/重建把背景粒子错误连进真链（False Positive）",
], Inches(0.9), Inches(1.5), Inches(11.6), Inches(5.0), size=18)

# ============ Slide 3: 问题诊断 ============
s = new_slide(3)
add_title_bar(s, "问题诊断：为什么信号链会丢失")
add_bullets(s, [
    "链存活 = 所有节点 AND 所有正边同时保留（级联 AND 条件）",
    "硬阈值剪枝 thr 0.2：edge 剪枝杀死 30.5% 的链，node 剪枝杀死 26.3%（高度重叠）",
    "Loss 分解：tt_edges(×33) 占 81-91%，node 仅 5.8%，LCA 仅 1.5% → node/LCA 梯度投入严重不足",
    "训练时全图学习、推理时硬剪枝 → train-inference gap",
    "方案 F（seed-expand）实测失败：9 种子→73 节点（占全图 80%）= 等于不剪 → NoneIso 96.9%",
    "教训：CERN 高连通图上“连通性”≈噪声本身，应关心“最物理”而非“最像树”",
], Inches(0.9), Inches(1.45), Inches(11.6), Inches(5.3), size=17)

# ============ Slide 4: 优化方案总览 ============
s = new_slide(4)
add_title_bar(s, "优化方案总览（5 项改进）")
rows = [
    ("B2 可微剪枝", "训练时软掩码模拟剪枝，消除 train-inference gap（核心）"),
    ("源检测头", "Rumor Centrality 训练化：监督 GNN 学“找根”（第6头）"),
    ("Loss 再平衡", "w_node=10, w_lca=10：node/LCA 梯度占比提到 20-30%"),
    ("LCA weights clip", "修复 CERN class2/3 零样本导致的 inf 权重"),
    ("链级 LCA 判据", "推理侧：链内边平均 LCA 置信度过滤噪声链"),
]
kv_rows(s, rows, top=1.7, row_h=0.95, k_w=2.4, size=17)

# ============ Slide 5: B2 可微剪枝 ============
s = new_slide(5)
add_title_bar(s, "B2：可微剪枝训练（核心，消除 gap）")
add_bullets(s, [
    "做法：训练时对 node/edge weight 施加可微软掩码 mask = σ((w − cut) / τ)",
    "消息传递在“被剪的图”上进行，梯度穿过掩码流回主干",
    "阈值式软掩码（弃用 top-k：方案 F 已证明 top-k 在 CERN 上失控）",
    "温度退火：τ 1.0 → 0.1（跨 100 epoch），从软到近硬",
    "仅在最后一个 GN block 启用（与推理剪枝位置对齐），eval 自动关闭",
    "参数：b2_cut=0.5, b2_k=0（纯阈值）, b2_tau_start=1.0, b2_tau_end=0.1",
], Inches(0.9), Inches(1.5), Inches(11.6), Inches(5.0), size=17)

# ============ Slide 6: 源检测头 ============
s = new_slide(6)
add_title_bar(s, "源检测头：Rumor Centrality 训练化")
add_bullets(s, [
    "动机：真实衰变链是有清晰“根”（B 介子）的树；噪声团块无根",
    "truth 链 → 根 = 链内 rumor centrality 最大的节点（truth 图结构算，无噪声）",
    "head：节点级二分类“该节点是否为某链的根”，BCE 损失",
    "RC 位于重建最后一步，面对相对干净的信息 → 其“找根”能力可作监督信号进训练",
    "让主干显式学“根-叶结构”，与推理侧 RC（找根）对齐",
    "w_source=5.0；标签生成复用 chain_center_score（truth_chain_roots）",
], Inches(0.9), Inches(1.5), Inches(11.6), Inches(5.0), size=17)

# ============ Slide 7: Loss 再平衡 + LCA clip ============
s = new_slide(7)
add_title_bar(s, "Loss 再平衡 + LCA weights clip")
add_bullets(s, [
    "原 combined = LCA + 1·t_nodes + 33·tt_edges + pv_asso",
    "tt_edges 占 91%，node/LCA 欠投入 → 模型动力学不足",
    "新 combined = LCA + 10·t_nodes + 33·tt_edges + pv_asso + 5·source_loss",
    "LCA weights clip：CERN 数据 class2/3 样本为 0 → 权重 inf → CrossEntropyLoss inf",
    "修复：nan_to_num(inf→1e3) + clamp(≤1e3)，与 FT 处理一致",
], Inches(0.9), Inches(1.5), Inches(11.6), Inches(5.0), size=18)

# ============ Slide 8: 链级 LCA 判据 ============
s = new_slide(8)
add_title_bar(s, "链级 LCA 物理判据（推理侧）")
add_bullets(s, [
    "“最物理”而非“最像树”：真链的链内边是模型高置信的物理关系",
    "conf = 链内边被判类别的平均 softmax 概率（高=物理自洽）",
    "过滤：conf < thr 的链剔除 → 压低 NoneIso",
    "与“放宽剪枝”组合：放宽剪枝（thr 0.7）管“别剪掉真的”，LCA 过滤管“别留下假的”",
    "评估时扫描 conf_thr ∈ {0.0, 0.3, 0.5}；chain_lca_record 模式可后处理多阈值",
], Inches(0.9), Inches(1.5), Inches(11.6), Inches(5.0), size=17)

# ============ Slide 9: v36 训练结果 ============
s = new_slide(9)
add_title_bar(s, "v36（B2+源检测头）训练结果")
add_bullets(s, [
    "起点 v35 best（ep68, val 35.955）→ v36 续训至 ep94（早停）",
    "val best 35.58 @ep79（优于 v35 的 35.70）",
    "B2 温度退火正常：ep70 τ=0.37 → ep80 τ=0.28",
    "source_loss 缓慢下降：0.0387 → 0.0338（源检测头在学习）",
    "⚠️ v36 自动评估用 thr0.7（口径与 v31 基线 thr0.9 不同）→ 需同口径对比",
    "已提交：v36+thr0.9（公平对比）、v36+thr0.7+chain_lca（组合验证）",
], Inches(0.9), Inches(1.5), Inches(11.6), Inches(5.0), size=17)

# ============ Slide 10: 评估矩阵（待完成） ============
s = new_slide(10)
add_title_bar(s, "评估矩阵（训练后）")
add_bullets(s, [
    "A. 基线（v31, thr 0.9）：Perfect 23.9%, NoneIso 56.6%  ← 对照",
    "B. 放宽剪枝（thr 0.7）：降低真链误杀",
    "C. B + chain_lca_filter（conf 0.3/0.5 扫描）：链级物理判据（主过滤）",
    "D. + B2 训练（v36）：打分更可信 → 过滤更有效",
    "E. D + 源检测头叠加验证",
    "指标：per-chain Perfect/NoneIso, per-event Perfect, 候选边数",
], Inches(0.9), Inches(1.5), Inches(11.6), Inches(5.0), size=17)

# ============ Slide 11: 下一步 ============
s = new_slide(11)
add_title_bar(s, "下一步")
add_bullets(s, [
    "等待两个评估作业完成（v36+thr0.9 公平对比 / v36+thr0.7+chain_lca）",
    "用同口径（thr0.9）对比 v31 vs v36，判断 B2+源检测头的真实收益",
    "v35 自动评估因 GPU 驱动崩溃（CUDA driver error），用独立评估脚本重跑",
    "若 B2 有效 → 继续调 b2_cut / w_source；若无效 → 回退到 thr0.9 + chain_lca 组合",
], Inches(0.9), Inches(1.5), Inches(11.6), Inches(5.0), size=18)

prs.save(OUT)
print("saved ->", OUT)
