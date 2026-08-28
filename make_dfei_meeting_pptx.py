#!/usr/bin/env python3
"""DFEI group meeting PPT (English): bug fixes -> optimizations -> physics supervision -> open questions.

结果页留 TBD 占位, 训练完成后更新 RESULTS 字典重新生成即可。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = '/lzufs/home/guoqingxiang/dfei/scalable_mtl_hgnn'
FIG = BASE + '/meeting_figs'
OUT = FIG + '/DFEI_group_meeting_EN.pptx'
os.makedirs(FIG, exist_ok=True)

BLUE = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT = RGBColor(0xEA, 0xF2, 0xF8)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
FONT = 'Calibri'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title_bar(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = BLUE
    r.font.name = FONT
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(13)
        r2.font.color.rgb = GRAY
        r2.font.name = FONT
    ln = slide.shapes.add_shape(1, Inches(0.6), Inches(1.05), Inches(12.1), Pt(3))
    ln.fill.solid()
    ln.fill.fore_color.rgb = BLUE
    ln.line.fill.background()


def add_bullets(slide, items, left, top, width, height, size=15, color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        r = p.add_run()
        r.text = ('    ' * lvl) + ('•  ' if lvl == 0 else '–  ') + txt
        r.font.size = Pt(size - 2 * lvl)
        r.font.color.rgb = color or DARK
        r.font.name = FONT
    return tb


def add_table(slide, data, left, top, width, height, col_widths=None, font_size=12, header_fill=BLUE):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = w
    for i in range(rows):
        for j in range(cols):
            cell = tbl.cell(i, j)
            cell.text = str(data[i][j])
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.name = FONT
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i == 0 else DARK
                p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    return tbl


def add_footer(slide, idx):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = 'Qingxiang Guo · UCAS · DFEI group meeting      ' + str(idx)
    r.font.size = Pt(10)
    r.font.color.rgb = GRAY
    r.font.name = FONT


def new_slide(idx):
    s = prs.slides.add_slide(BLANK)
    add_footer(s, idx)
    return s


# ============ RESULTS: 训练完成后更新这些值 ============
# 正在跑的实验 (TBD): 10014718 mass+struct+mom 组合; 9995727 公开数据 v38; 10014638 masshead2 best 评估
R_TRAIN_RUNNING = "TBD (training in progress)"
R_PUBLIC = "TBD (training in progress, ~27/100 ep, val_combined 42.3 best so far)"
R_BEST_EVAL = "TBD (queued)"

# ============ S1 Title ============
s = new_slide(1)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = 'Optimizing DFEI: Bug Fixes, Physics-Supervised Learning, and Open Questions'
r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = FONT
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = 'A work-in-progress report built on the DFEI prototype (García Pardiñas et al., arXiv:2304.08610)'
r.font.size = Pt(18); r.font.color.rgb = DARK; r.font.name = FONT
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.9))
tf = tb.text_frame
r = tf.paragraphs[0].add_run()
r.text = 'Qingxiang Guo · University of Chinese Academy of Sciences · 2026-08'
r.font.size = Pt(14); r.font.color.rgb = GRAY; r.font.name = FONT

# ============ S2 Overview ============
s = new_slide(2)
add_title_bar(s, 'Overview')
add_bullets(s, [
    'Part 1 — Starting point: state of the DFEI code when I took it over',
    'Part 2 — Bug fix and the optimization line (v31 → v38)',
    'Part 3 — New direction: output-side physics supervision (mass / structure / momentum heads)',
    ('with PhyIP-style linear-probe evidence', 1),
    'Part 4 — Results so far (some still training)',
    'Part 5 — Open questions for the group: what is the right application form?',
    ('trigger-line assistance · flavor tagging · lightweight deployment', 1),
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.4), size=17)

# ============ S3 Starting point ============
s = new_slide(3)
add_title_bar(s, 'Starting Point: What I Found in the Repo', 'Baseline = the public prototype trained on MC, v31-era code')
add_bullets(s, [
    'DFEI = full-event GNN: simultaneously classify, isolate and hierarchically reconstruct all heavy-hadron decay chains per event',
    'Four supervised tasks: LCAG edge classification (4 classes), node pruning, edge pruning, PV association',
    'Problems I found in the inherited code:',
    ('Class-weight bug: config key LCA__weights (double underscore) silently disabled the weights', 1),
    ('→ class1 (parent–child) accuracy stuck at ~0%: the model only learned the dominant class0 background', 1),
    ('Hard-threshold pruning at inference vs full-graph training → train-inference gap', 1),
    'Baseline numbers (20 test files, thr 0.9): Perfect ≈ 23.9%',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S4 Bug fix ============
s = new_slide(4)
add_title_bar(s, 'Bug Fix + First Win: Inverse-Frequency Class Weighting', 'Lesson: make sure the supervision signal actually reaches the loss')
add_bullets(s, [
    'Fix: LCA__weights → LCA_weights, weights = N / (4·n_c) per class (clamped)',
    'Result: class1 accuracy recovered from ~0% to >70%; Perfect 23.9% → 26.3% (v31 → v36)',
    'Takeaway: verify the signal enters the loss before tuning the model — this unlocked the whole line',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(3.2), size=16)

# ============ S5 B2 pruning ============
s = new_slide(5)
add_title_bar(s, 'B2: Differentiable Pruning + Temperature Annealing', 'Train on the same (pruned) graphs the inference sees')
add_bullets(s, [
    'Gap: training runs on the full graph; inference prunes nodes/edges first (hard thresholds)',
    'Soft mask at the last GN block: w_eff = w · σ((w − cut)/τ), τ annealed 1.0 → 0.1',
    ('τ large → smooth mask (easy gradients); τ small → approximates the hard 0.9 threshold', 1),
    'cut aligned to inference: 0.5 → 0.7 → 0.85',
    'Effect: model learns to push background edges down / signal edges up explicitly',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.6), size=15)

# ============ S6 class2 + chain_lca ============
s = new_slide(6)
add_title_bar(s, 'class2 Weighting + In-chain LCA Consistency', 'Sister (same-mother) edges are the structural bottleneck')
add_bullets(s, [
    'class2 (same B) accuracy was the weak point for chain structure; class2 CE weight 3.0 → 2.0 (3.0 hurt class1)',
    'In-chain consistency losses (truth-chain edges only):',
    ('Hinge: penalize low-confidence chain edges (margin 0.3) → chains become cleaner at inference', 1),
    ('CE on chain-edge classes (chain_lca_ce): directly supervise class1/2/3 on the ~0.1% structural edges', 1),
    'Net effect: v38 Perfect 26.3% → 29.3%',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.4), size=15)

# ============ S7 Source head ============
s = new_slide(7)
add_title_bar(s, 'Source Head: Rumor-Centrality Training', 'Teach the backbone the root–leaf structure of a chain')
add_bullets(s, [
    'Inference finds chain roots with Rumor Centrality (RC); training never saw roots',
    'Added a node head predicting "is this the chain root?" (label = RC-argmax node in each truth chain, BCE)',
    'Note for the group: RC-argmax is the centroid of the *track graph*, not necessarily the B itself',
    ('most visible tracks are final-state particles — parent–child *track* relations are physically rare', 1),
    ('this observation motivates our structure-supervision discussion later', 1),
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.6), size=15)

# ============ S8 Physics supervision ============
s = new_slide(8)
add_title_bar(s, 'New Direction: Output-Side Physics Supervision', 'Use physical quantities as auxiliary supervision, keep the end-to-end objective')
add_bullets(s, [
    'Idea: instead of hand-crafting physics features into the input, supervise the representations with physical targets',
    'Mass head (edge-level): regress log10(m_ππ) of each track pair from the edge representation',
    ('target computed from track momenta only — no extra labels, data-intrinsic', 1),
    ('m_ππ under the pion hypothesis; same-mother pairs sit at resonance masses → structural information', 1),
    'Effect (19 epochs, masshead2): AllParticles 52.1 → 55.9, Perfect 29.3 → 32.7, class2 44.7 → 51.1',
    'Also implementing: structure head (depth + RC regression) and momentum head (node-level p regression)',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.4), size=15)

# ============ S9 PhyIP evidence ============
s = new_slide(9)
add_title_bar(s, 'Evidence: PhyIP-Style Linear Probe', 'Does supervision make the representation physically readable?')
add_table(s, [
    ['Probe (frozen backbone, Ridge regression)', 'v38', 'masshead2'],
    ['Edge repr → log10(m_ππ)', 'R² = 0.003', 'R² = 0.930'],
    ['Node repr → p (px/py/pz)', 'R² ≈ 0', 'R² ≈ 0'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(1.8), font_size=14)
add_bullets(s, [
    'Mass supervision makes the edge representation linearly readable for the mass — the "physics is in the representation" claim (HEP-JEPA style) is confirmed',
    'Nodes are NOT linearly readable for momenta (graph_norm + ReLU scramble them) → momentum head (mom head) added to fix this',
], Inches(0.8), Inches(3.7), Inches(11.8), Inches(3.0), size=15)

# ============ S10 Head zoo ============
s = new_slide(10)
add_title_bar(s, 'The Head Zoo (9 supervision heads)', 'All train jointly with the shared backbone')
add_table(s, [
    ['Head', 'Target', 'Status'],
    ['LCA / node / edge / PV (core 4)', 'LCAG class, signal node/edge, PV assoc.', 'baseline'],
    ['chain_scorer (5)', 'candidate-chain likelihood', 'implemented, not yet trained'],
    ['source (6)', 'chain-root binary', 'trained (v36+)'],
    ['mass (7)', 'edge log10(m_ππ)', 'trained → +3.4pp Perfect'],
    ['struct (8)', 'node depth + RC', 'training (in progress)'],
    ['mom (9)', 'node normalized p regression', 'training (in progress)'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(3.6), font_size=13)

# ============ S11 Results summary (TBD) ============
s = new_slide(11)
add_title_bar(s, 'Results Summary (some still running)', 'Same 20 test files, thr 0.9')
add_table(s, [
    ['Model', 'AllParticles', 'Perfect', 'class2 acc'],
    ['v38 (baseline)', '52.1%', '29.3%', '~45%'],
    ['masshead2 (19ep)', '55.9%', '32.7%', '51.1%'],
    ['mass+struct+mom (running)', R_TRAIN_RUNNING, R_TRAIN_RUNNING, R_TRAIN_RUNNING],
    ['v38 best-ckpt eval (queued)', R_BEST_EVAL, R_BEST_EVAL, R_BEST_EVAL],
    ['Public-data v38 (running)', R_PUBLIC, R_PUBLIC, R_PUBLIC],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(2.8), font_size=13)
add_bullets(s, [
    'To be filled when trainings finish — this slide is a placeholder',
], Inches(0.8), Inches(4.6), Inches(11.8), Inches(1.0), size=13)

# ============ S12 Trigger perspective ============
s = new_slide(12)
add_title_bar(s, 'Trigger Perspective: Candidate-Chain Scoring', 'Are the outputs usable to assist HLT2-style selection?')
add_bullets(s, [
    'Question: full-event reconstruction (AllParticles→90%) is the wrong metric for a trigger line — what matters is candidate-chain classification',
    'Preliminary AUC (chain criteria from model outputs, 20 events):',
    ('truth chains vs model-pruned components (realistic background): AUC ≈ 0.78', 1),
    ('with easy random combos: AUC ≈ 0.90; strongest single feature: in-chain non-class0 probability 0.96', 1),
    'Not yet trained: chain_scorer (head 5) — trained scorer + event-context features should push this higher',
    'Time budget: current model ~12 ms/event on GPU; needs 10-100x compression for HLT2 (distillation, pruning)',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S13 Open questions ============
s = new_slide(13)
add_title_bar(s, 'Open Questions for the Group', 'We would like your input on the application form')
add_bullets(s, [
    'Q1 — Application form: should we aim at trigger-line *assistance* (candidate scoring) rather than full-event reconstruction? Which metric would the collaboration trust?',
    'Q2 — Track-level parent–child relations are physically rare (most tracks are final-state). Is "same-source clustering" (class2) the real problem DFEI should solve?',
    'Q3 — Invariant mass is computed analytically in HLT2, not by the network. What is the genuine added value of the GNN — structure/context, not kinematics?',
    'Q4 — Flavor tagging: a transformer-based tagger group reported good tagging power. Could DFEI event representations (same-source clusters) combine with it, e.g. for same-side tagging inputs?',
    'Q5 — Is there a realistic HLT2 (or Upgrade II) time budget for a lightweight GNN per event, or is offline/Turbo the right target?',
    'Q6 — Physics supervision (mass/depth/momentum heads) + linear-probe verification: is this a direction the group wants to pursue, and are there better physical targets?',
], Inches(0.8), Inches(1.3), Inches(11.8), Inches(6.0), size=14)

# ============ S14 Future ideas ============
s = new_slide(14)
add_title_bar(s, 'Some Ideas We Are Considering', 'Feedback welcome')
add_bullets(s, [
    'Train the chain scorer (head 5) and measure candidate-chain AUC properly (full 20-file evaluation)',
    'Distill the 4-block GNN to a small model and measure per-event CPU latency → quantify HLT2 feasibility',
    'Verify "efficiency independent of luminosity" claim on high-multiplicity (μ~50) MC — relevant for Upgrade II',
    'Combine DFEI same-source clustering with a transformer tagger for same-side flavor tagging',
    'Release the physics-supervision recipe (mass/struct/mom heads + linear probe) as a reproducible methodology',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.4), size=15)

# ============ S15 Backup ============
s = new_slide(15)
add_title_bar(s, 'Backup: Data & Setup')
add_bullets(s, [
    'Data: LHCb Upgrade-I conditions (Run 3), ~150 tracks/event, MC normalized, 200 train / 20 val / 20 test files',
    'Public DFEI dataset (converted_LHCbcollision) used for the independent v38 verification run (in progress)',
    'Hardware: single GPU (10-44 GB), batch 8, ~40 min/epoch',
    'Code: scalable_mtl_hgnn fork, all changes reproducible from config files',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.4), size=15)

prs.save(OUT)
print('[ok] 已保存: ' + OUT)
print('     共 ' + str(len(prs.slides._sldIdLst)) + ' 页')
