#!/usr/bin/env python3
"""生成中文详细 PPT: DFEI 原始架构详解 (到每个 MLP) + Bug 修复后每次升级的算法详解。
输出: architecture_figs/DFEI_architecture_algorithms.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = '/lzufs/home/guoqingxiang/dfei/scalable_mtl_hgnn'
FIG = BASE + '/architecture_figs'
OUT = FIG + '/DFEI_architecture_algorithms.pptx'
os.makedirs(FIG, exist_ok=True)

BLUE = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT = RGBColor(0xEA, 0xF2, 0xF8)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
FONT = 'Microsoft YaHei'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title_bar(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = BLUE
    r.font.name = FONT
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(13)
        r2.font.color.rgb = GRAY
        r2.font.name = FONT
    ln = slide.shapes.add_shape(1, Inches(0.6), Inches(1.05), Inches(12.1), Pt(3))
    ln.fill.solid()
    ln.fill.fore_color.rgb = BLUE
    ln.line.fill.background()


def add_bullets(slide, items, left, top, width, height, size=15, color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        r = p.add_run()
        r.text = ('    ' * lvl) + ('•  ' if lvl == 0 else '–  ') + txt
        r.font.size = Pt(size - 2 * lvl)
        r.font.color.rgb = color or DARK
        r.font.name = FONT
    return tb


def add_table(slide, data, left, top, width, height, col_widths=None, font_size=12, header_fill=BLUE):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = w
    for i in range(rows):
        for j in range(cols):
            cell = tbl.cell(i, j)
            cell.text = str(data[i][j])
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.name = FONT
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i == 0 else DARK
                p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    return tbl


def add_footer(slide, idx):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = 'Qingxiang Guo · UCAS · DFEI 架构与算法手册      ' + str(idx)
    r.font.size = Pt(10)
    r.font.color.rgb = GRAY
    r.font.name = FONT


def new_slide(idx):
    s = prs.slides.add_slide(BLANK)
    add_footer(s, idx)
    return s


def add_box(slide, left, top, width, height, text, fill=BLUE, font_color=RGBColor(0xFF, 0xFF, 0xFF), size=13, bold=True):
    sh = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = font_color
    r.font.name = FONT
    return sh


def add_arrow(slide, x1, y1, x2, y2):
    ln = slide.shapes.add_shape(8, Inches(x1), Inches(y1), Inches(max(x2 - x1, 0.05)), Pt(3))
    ln.fill.solid()
    ln.fill.fore_color.rgb = GRAY
    ln.line.fill.background()


# ============ S1 封面 ============
s = new_slide(1)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.4))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = 'DFEI 架构详解与优化算法手册'
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = FONT
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = '从原始异构 GNN 架构（具体到每个 MLP），到每次升级尝试的具体算法'
r.font.size = Pt(20); r.font.color.rgb = DARK; r.font.name = FONT
r = tf.add_paragraph().add_run()
r.text = 'LCAG 分类 · 节点/边剪枝 · PV 关联 · B2 可微剪枝 · 温度退火 · Gumbel-Softmax · 课程学习'
r.font.size = Pt(15); r.font.color.rgb = GRAY; r.font.name = FONT
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.8))
tf = tb.text_frame
r = tf.paragraphs[0].add_run()
r.text = 'Qingxiang Guo · 中国科学院大学 · 2026-08'
r.font.size = Pt(15); r.font.color.rgb = GRAY; r.font.name = FONT

# ============ S2 目录 ============
s = new_slide(2)
add_title_bar(s, '目录')
add_bullets(s, [
    '第一部分：DFEI 原始架构详解（具体到每个 MLP 的功能）',
    ('背景：事件 → 图，多任务全景', 1),
    ('数据流：Encoder → GN Block ×4 → Decoder → 输出头', 1),
    ('四个监督任务：LCAG 分类 / 节点剪枝 / 边剪枝 / PV 关联', 1),
    '第二部分：Bug 修复后每次升级的具体算法',
    ('v31 class weight 修复（逆向频率加权）', 1),
    ('v36 B2 可微剪枝 + 温度退火、源检测头（Rumor Centrality）', 1),
    ('v37/38 class2 加权、链内 LCA 一致性（hinge + CE）', 1),
    ('v40-42 PV 分簇：Gumbel-Softmax、课程学习、val 对齐、子图 cap', 1),
    '第三部分：经验总结与当前状态',
    '第四部分：最新进展（2026-08-25~26）',
    ('方向调整：评估口径转向 trigger（AllParticles 优先）', 1),
    ('输出侧物理监督：边级不变质量回归（mass head）', 1),
    ('EarlyStopping 续训 bug 修复', 1),
    ('PV 分簇（DA + 重叠软成员）评估结论', 1),
    ('公开数据验证 v38（进行中）', 1),
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=17)

# ============ S3 背景 ============
s = new_slide(3)
add_title_bar(s, '背景：DFEI 要解决什么问题', 'Deep Full Event Interpretation — LHCb B 介子衰变链重建')
add_bullets(s, [
    'LHCb 对撞事件：一次质子-质子对撞产生大量径迹（track）和顶点（PV，Primary Vertex）',
    '物理目标：找出哪些径迹属于同一条 B 介子衰变链（B → 中间粒子 → 末态粒子）',
    '挑战：',
    ('事件里绝大多数径迹是背景（来自同一 PV 的其他碰撞/喷注），信号链只占很小比例', 1),
    ('一个事件有多个 PV（堆叠），信号 B 可能来自其中某一个', 1),
    ('衰变链有多级：B → J/ψ → μμ，需要判断径迹之间的父子/同母关系', 1),
    'DFEI 的思路：把事件建模成图，用异构 GNN 同时学会：',
    ('① 哪些径迹是信号（节点剪枝）  ② 径迹之间是否在一条链上、是什么关系（边剪枝 + LCAG 分类）', 1),
    ('③ 径迹属于哪个 PV（PV 关联）', 1),
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=16)

# ============ S4 事件→图 ============
s = new_slide(4)
add_title_bar(s, '事件 → 图：DFEI 的输入表示')
add_bullets(s, [
    '节点（node_types）：',
    ('tracks 节点：每条径迹一个，特征 x（运动学：动量、角度、IP 等，8 维）+ PID 信息（拼接后 16 维）', 1),
    ('pvs 节点：每个主顶点一个，特征 x（3 维，顶点位置相关）', 1),
    ('globals 节点：事件级全局信息（每事件 1 个，2 维）', 1),
    '边（edge_types）：',
    ('tracks → tracks 边：径迹两两之间，特征 edges（5 维物理量，如两径迹角度差、顶点距离），标签 y = LCA 类别', 1),
    ('tracks → pvs 边：径迹与每个 PV 之间（稠密，每 track 连所有 PV），特征 edges（1 维，IP 类），标签 y = 是否属于该 PV', 1),
    '标签（监督信号）来自 MC 真值：',
    ('truth 衰变树 → 每条 tt 边的 LCAG 类别（0-3）', 1),
    ('truth 径迹/边是否信号 → 节点/边剪枝标签', 1),
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=16)

# ============ S5 多任务全景 ============
s = new_slide(5)
add_title_bar(s, '多任务全景：模型同时预测四件事')
add_table(s, [
    ['任务', '预测什么', '监督标签', '损失', '用于'],
    ['① LCAG 分类', '每条 tt 边的 LCA 类别（4 类）', 'truth LCA 类别 y∈{0,1,2,3}', 'CrossEntropy（加权）', '链重建（重建衰变树）'],
    ['② 节点剪枝', '每条径迹是信号的概率（0~1）', 'ft≠1（信号）', 'BCE（加权）', '砍掉背景径迹'],
    ['③ 边剪枝', '每条 tt 边是链内边的概率（0~1）', 'LCA y>0（链内）', 'BCE（加权）', '砍掉背景边'],
    ['④ PV 关联', '径迹属于每个 PV 的概率（0~1）', 'track-PV 关联真值', 'BCE（加权）', '径迹-顶点归属'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(3.4), font_size=13)
add_bullets(s, [
    'LCAG 四类的含义（两径迹在 truth 衰变树里的最低共同祖先关系）：',
    ('class 0：无共同祖先（背景对）   class 1：直系母子（一个直接衰变成另一个）', 1),
    ('class 2：同母（姊妹，同一中间粒子衰变产物）   class 3：祖孙（隔 ≥2 代）', 1),
], Inches(0.8), Inches(5.2), Inches(11.8), Inches(1.8), size=14)

# ============ S6 整体数据流 ============
s = new_slide(6)
add_title_bar(s, 'DFEI 整体数据流（model.py: DFEI_HGNN.forward）', '一个异构 GNN：Encoder → GN Block ×4 → Decoder → 输出头')
add_box(s, 0.9, 1.6, 1.7, 1.1, '输入图\n(原始特征)', RGBColor(0x8E, 0x8E, 0x8E))
add_box(s, 2.9, 1.6, 1.7, 1.1, 'Encoder\n(编码MLP)', BLUE)
add_box(s, 4.9, 1.6, 2.6, 1.1, 'GN Block ×4\n(消息传递)', BLUE)
add_box(s, 7.8, 1.6, 1.7, 1.1, 'Decoder\n(解码MLP)', BLUE)
add_box(s, 9.9, 1.6, 2.6, 1.1, '输出头\n(4个任务)', ORANGE)
add_arrow(s, 2.6, 1.9, 2.9, 1.9)
add_arrow(s, 4.6, 1.9, 4.9, 1.9)
add_arrow(s, 7.5, 1.9, 7.8, 1.9)
add_arrow(s, 9.5, 1.9, 9.9, 1.9)
add_bullets(s, [
    '关键细节：GN Block 之间用残差拼接——每过一个 block，把特征与潜变量图（encoder 输出）按最后一维拼接（hetero_graph_concat），避免深层退化',
    'Encoder/Decoder 每个元素类型（节点/边/全局）各有一个独立 MLP → 异构（不同对象用不同网络）',
    'B2 可微剪枝只在最后一个 GN Block 启用（与推理剪枝位置对齐），其余 Block 正常学全图表征',
], Inches(0.9), Inches(3.2), Inches(11.6), Inches(3.4), size=15)

# ============ S7 Encoder ============
s = new_slide(7)
add_title_bar(s, 'Encoder：每个元素类型一个编码 MLP（hetero_graphcoder.py）')
add_table(s, [
    ['MLP', '输入', '输出维度', '作用'],
    ['tracks 编码 MLP', 'tracks.x（8 维原始运动学，use_pid 时拼接 PID → 16）', '16', '把径迹原始特征映射到隐空间'],
    ['pvs 编码 MLP', 'pvs.x（3 维）', '16', 'PV 特征 → 隐空间'],
    ['global 编码 MLP', 'globals.x（2 维）', '16', '事件全局特征 → 隐空间'],
    ['tracks_tracks 边编码 MLP', 'tt 边 edges（5 维物理量）', '16', '径迹对特征 → 隐空间'],
    ['tracks_pvs 边编码 MLP', 'tr-pv 边 edges（1 维 IP 类）', '16', '径迹-PV 关联特征 → 隐空间'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(3.0), font_size=13)
add_bullets(s, [
    '结构：MLP(layers=[输入,128,128,16])，即 3 个线性层，每层 ReLU + 可选 norm（graph_norm）+ dropout 0.01',
    '输入维度用 -1 占位，构建时根据实际数据自动推断',
    '作用一句话：把异构的原始特征统一编码进同一 16 维隐空间，供后续消息传递使用',
], Inches(0.8), Inches(4.8), Inches(11.8), Inches(2.0), size=15)

# ============ S8 GN Block 三阶段 ============
s = new_slide(8)
add_title_bar(s, 'GN Block 内部：消息传递三阶段（hetero_graph_network.py）', '每个 Block = Edge block → Node block → Global block，各有一个 MLP')
add_box(s, 0.9, 1.7, 3.6, 1.0, '① Edge block MLP\n更新每条边的特征', BLUE)
add_box(s, 4.9, 1.7, 3.6, 1.0, '② Node block MLP\n更新每个节点的特征', BLUE)
add_box(s, 8.9, 1.7, 3.6, 1.0, '③ Global block MLP\n更新事件全局特征', BLUE)
add_arrow(s, 4.5, 2.0, 4.9, 2.0)
add_arrow(s, 8.5, 2.0, 8.9, 2.0)
add_bullets(s, [
    '这是 Graph Network（GN）的边→节点→全局更新范式：信息先在边上更新，再聚合到节点，最后聚合到全局',
    '每次消息传递后，边/节点特征都携带了邻域信息 → 多层之后感受野覆盖整个事件',
    '与 B2 的配合：② 里消息传递用 edge_weights 加权（weighted_pass），B2 软掩码就作用在这些权重上',
], Inches(0.9), Inches(3.1), Inches(11.6), Inches(3.0), size=15)

# ============ S9 Edge block ============
s = new_slide(9)
add_title_bar(s, '① Edge block MLP：更新边特征（hetero_edge_block.py）')
add_table(s, [
    ['拼接的输入（concat）', '含义'],
    ['edges（当前边特征 16 维）', '这条边自己的信息'],
    ['接收节点 x（edge_index[1]）', '边终点（receiver）的节点特征'],
    ['发送节点 x（edge_index[0]）', '边起点（sender）的节点特征'],
    ['全局 x（按起点 batch 广播）', '事件级信息'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(3.0), font_size=13)
add_bullets(s, [
    'MLP(layers=[-1,128,128,16])：拼接后 → 3 层线性 + ReLU → 更新后的边特征（仍 16 维）',
    '作用：让每条边看到自己两个端点的状态和事件全局，从而更新自己的表示——边特征从此携带结构信息',
], Inches(0.8), Inches(4.8), Inches(11.8), Inches(2.0), size=15)

# ============ S10 Node block ============
s = new_slide(10)
add_title_bar(s, '② Node block MLP：更新节点特征（hetero_node_block.py）')
add_table(s, [
    ['拼接的输入（concat）', '含义'],
    ['聚合的发送边消息（HeteroEdgesToNodesAggregator）', '把从该节点出发的所有边特征按边权重加权求和 → 邻居信息'],
    ['聚合的接收边消息（receiver）', '指向该节点的边特征加权求和'],
    ['节点自身 x', '节点原来的特征'],
    ['全局 x（按节点 batch 广播）', '事件级信息'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(3.0), font_size=13)
add_bullets(s, [
    '聚合器：对每个节点，把关联边的特征（16 维）按对应边权重 edge_weights 加权求和 → 相当于邻居消息汇聚（消息传递）',
    'MLP(layers=[-1,128,128,16])：拼接后更新节点特征（仍 16 维）',
    '作用：每个节点融合了邻居边的信息 → 径迹知道和它相连的径迹是谁、关系如何',
], Inches(0.8), Inches(4.8), Inches(11.8), Inches(2.0), size=15)

# ============ S11 Global + Decoder + 输出头 ============
s = new_slide(11)
add_title_bar(s, '③ Global block MLP + Decoder + 输出头')
add_bullets(s, [
    '③ Global block MLP（hetero_global_block.py）：',
    ('输入 = 所有边特征聚合（加权求和）+ 所有节点特征聚合 + 全局自身 → MLP 更新全局特征（16 维）', 1),
    ('作用：把整张图的信息压缩进事件级表示，再广播回所有边/节点（下一层的 context）', 1),
    'Decoder（HeteroGraphCoder，与 Encoder 同结构）：把隐特征解码回更利于分类的空间',
    '输出头（hetero_output_trafo.py / hetero_graph_network.py 里的 MLP_infer）：',
    ('LCA 头：tt 边 → MLP → 4 维 logits（CrossEntropy）', 1),
    ('MLP_infer（节点）：tracks 节点 → MLP → 1 维 logit → sigmoid = 节点保留概率（节点剪枝）', 1),
    ('MLP_infer（边）：每种边 → 同一 MLP → 1 维 logit → sigmoid = 边保留概率（边剪枝 / PV 关联）', 1),
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S12 推理重建 ============
s = new_slide(12)
add_title_bar(s, '推理重建流程（reconstruction.py）', '训练学到的权重，在推理时怎么用')
add_bullets(s, [
    '① 节点剪枝：node_weight > 0.9 的径迹保留（默认阈值 0.9，训练里 B2 软掩码模拟的就是这一步）',
    '② 边剪枝：edge_weight > 0.9 的 tt 边保留',
    '③ 得到剪枝后的稀疏图，对每条边取 LCA logits 的 argmax → 4 类判决',
    '④ 链重建（reconstruct_decay）：从剪枝后的图出发，按 LCA 关系把边串成衰变链',
    ('贪心：把母子关系的边连起来，形成一条条从 B 到末态粒子的链', 1),
    ('若一条边被判为 class2（同母）而连接关系冲突，会造成链分裂/错连 → 这正是 class2 准确率的痛点', 1),
    '★ 关键 Gap：训练时模型看全图，推理时先剪枝再重建——训练和推理的图结构不一致，这就是 B2 可微剪枝要解决的问题',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S13 优化总览 ============
s = new_slide(13)
add_title_bar(s, '优化总览（bug 修复后每一轮）')
add_table(s, [
    ['版本', '改动', '核心算法'],
    ['v31', '修复 class weight bug（LCA__weights→LCA_weights）', '逆向频率加权'],
    ['v36', 'B2 可微剪枝 + 源检测头', 'Sigmoid 软掩码 + 温度退火；Rumor Centrality'],
    ['v37', 'class2 专项加权 + 链内 LCA 一致性', '损失加权；Hinge loss'],
    ['v38', 'b2_cut 0.7→0.85 + chain_lca_ce', '阈值对齐；链内边类别 CE'],
    ['v39-42', 'PV 分簇分层重建（尝试）', 'Gumbel-Softmax、温度退火、课程学习、val 对齐'],
    ['v38+masshead', '第7个监督头：边级不变质量回归（输出侧物理监督）', 'SmoothL1 log10(ππ 不变质量)'],
], Inches(0.8), Inches(1.6), Inches(11.8), Inches(3.4), font_size=14)
add_bullets(s, [
    '结果里程碑（thr0.9 Perfect 比例）：v31 23.93% → v36 26.28% → v37 27.25% → v38 29.26%',
], Inches(0.8), Inches(5.2), Inches(11.8), Inches(1.2), size=15)

# ============ S14 v31 逆向频率加权 ============
s = new_slide(14)
add_title_bar(s, 'v31：class weight bug 修复 = 逆向频率加权（Inverse Frequency Weighting）', '类别不平衡问题的标准解法')
add_bullets(s, [
    '问题：LCAG 类别极度不平衡——class0（无关边）占绝对多数，class1/2/3（结构边）极少（约 0.1%）',
    '后果：模型只要全判 class0 就能把平均损失压得很低 → 结构边完全不学 → class1 准确率掉到 ~0%',
    '解法：给每个类一个权重，样本越少的类权重越大：',
    ('w_c = N / (K × n_c)，其中 N=总样本数，K=类别数，n_c=类 c 的样本数（代码：sum/(4*count)，再 clamp 到 ≤1e3）', 1),
    ('训练时 loss = Σ w_c × CE(预测, 类c) → 少数类（class1/2/3）被放大，模型被迫学结构边', 1),
    '配套 bug：配置键拼写错误 LCA__weights（双下划线）→ 权重根本没生效，修成 LCA_weights 后 class1 才被学起来',
    '⭐ 启发：先确认信号有没有进到损失里，再谈调模型——这是整个优化线的起点',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S15 B2 为什么 ============
s = new_slide(15)
add_title_bar(s, 'v36：B2 可微剪枝 ① 为什么需要它', 'Train-Inference Gap：训练看全图，推理先剪枝')
add_bullets(s, [
    '推理侧：先用阈值把节点/边剪掉（硬剪枝），再在稀疏图上做链重建',
    '训练侧（v31-35）：模型在全图上训练，从没见过被剪过的图',
    '后果：训练时的特征分布 ≠ 推理时的特征分布 → 剪枝后的结构预测不准',
    'B2 的思路：在训练的最后一步，用软掩码模拟推理的硬剪枝，让训练和推理看同一类图',
    ('关键设计：只在最后一个 GN Block 启用——因为推理时剪枝正是作用于最终输出的权重', 1),
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=16)

# ============ S16 B2 软掩码 ============
s = new_slide(16)
add_title_bar(s, 'v36：B2 可微剪枝 ② Sigmoid 软掩码', '硬剪枝不可导，用连续掩码代替')
add_bullets(s, [
    '硬剪枝：weight > 0.9 保留，否则丢弃 —— 不可导，无法反向传播',
    '软掩码（可导的近似）：',
    ('mask = σ((w − cut) / τ)，σ 是 sigmoid，cut 是剪枝阈值（b2_cut，默认 0.85），τ 是温度', 1),
    ('训练时把权重乘上掩码：w_effective = w × σ((w − cut)/τ)', 1),
    ('效果：w 远高于 cut → mask≈1（保留）；w 远低于 cut → mask≈0（近似剪掉）；w≈cut → 平滑过渡', 1),
    ('梯度可以流经 mask → 模型学会把该剪的边压低、该留的边抬高', 1),
    '阈值对齐：b2_cut 从 0.5 → 0.7 → 0.85 逐步对齐推理阈值 0.9，让模拟更接近真实剪枝',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=16)

# ============ S17 B2 温度退火 ============
s = new_slide(17)
add_title_bar(s, 'v36：B2 可微剪枝 ③ 温度退火（Temperature Annealing）', '为什么掩码要从软变硬')
add_bullets(s, [
    '问题：σ((w−cut)/τ) 里温度 τ 控制掩码的陡峭程度',
    ('τ 大（如 1.0）→ sigmoid 平滑，掩码在 cut 附近过渡缓和（软，留有余地，利于梯度流动）', 1),
    ('τ 小（如 0.1）→ sigmoid 陡峭，掩码趋近 0/1 硬判决（硬，与推理阈值剪枝一致）', 1),
    '退火策略：训练初期 τ=1.0（软，模型先学到大方向），随 epoch 线性降到 0.1（硬，模拟真实剪枝）',
    ('τ(epoch) = 1.0 + (0.1 − 1.0) × min(epoch/100, 1)', 1),
    '类比：模拟退火——先允许模糊，再逐步精确，避免一开始就被硬判决卡死梯度',
    '⭐ 通用套路：凡是离散决策需要放进训练，就用连续松弛（sigmoid/softmax）+ 温度退火',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=16)

# ============ S18 源检测头 ============
s = new_slide(18)
add_title_bar(s, 'v36：源检测头 = Rumor Centrality 训练化', '让 GNN 学会一条链的根在哪')
add_bullets(s, [
    '背景：推理时重建链需要找到根（B 介子候选），Rumor Centrality（RC）是一种在树上找根的方法',
    'Rumor Centrality 直觉：给定一棵树（链），信息传播从某个源头扩散；RC 给每个节点打一个它是源头的分数，最大者即最可能的根',
    '问题：推理时才用 RC 找根，训练时模型从没学过根-叶结构',
    '做法：新增第 6 个监督头 source_head（MLP：节点特征+节点权重 → 根概率 logit）',
    ('标签：对每条 truth 链，取 RC 分数最大的节点作为根标签（真值，无噪声）', 1),
    ('损失：BCE(logit, 根标签)，与主干联合训练 → 让主干显式学链的层级结构', 1),
    '效果：训练时找根与推理时 RC 找根对齐，主干对链结构的理解更深',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S19 class2 加权 ============
s = new_slide(19)
add_title_bar(s, 'v37：class2 专项加权', '同 B 边（姊妹关系）是链结构的关键短板')
add_bullets(s, [
    '观察：v36 里 class2（同母边）准确率 ~36%，是 Perfect 率的最大结构瓶颈',
    '原因：class1（母子）保证链的连通，class2（姊妹）决定同一个中间粒子下的多条子链怎么并成一条，判错直接破坏链结构',
    '做法：在逆向频率权重基础上再乘一个系数：lca_class2_weight = 3.0（v37）→ 2.0（v38）',
    ('lca_w[2] = 原权重 × lca_class2_weight → class2 的 CE 损失被进一步放大', 1),
    'v37 用 3.0 后发现 class1 反而下降（加权过度，模型顾此失彼）→ v38 降到 2.0',
    '⭐ 启发：专项加权要点到为止，权重过大破坏其他类别的学习',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=16)

# ============ S20 chain_lca ============
s = new_slide(20)
add_title_bar(s, 'v37/38：链内 LCA 一致性（chain_lca_loss + chain_lca_ce）', '把最物理的判据直接放进训练')
add_bullets(s, [
    '物理直觉：真链的链内边（y>0）应该是模型高置信的非背景边——这是推理侧 chain_lca_filter 的判据',
    '做法①（chain_lca_loss，Hinge）：只对 truth 链内边施加：',
    ('对链内边取被判类别的 softmax 概率 conf，loss = max(0, margin − conf)²，margin=0.3', 1),
    ('conf < margin 的边受罚 → 逼迫链内边高置信，让链在推理时更干净', 1),
    ('Hinge loss 特点：超过 margin 就不再罚（不像 CE 无上限），只关注不够自信的样本', 1),
    '做法②（chain_lca_ce，v38 升级）：链内边额外加类别 CE',
    ('直接用真类别（1/2/3）监督链内边的分类，对抗 class0 对主 LCA loss 的稀释', 1),
    '本质：把推理后处理（chain_lca_filter 阈值过滤）的判据前移到训练，让模型主动产出物理自洽的链',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S21 PV 分簇动机 ============
s = new_slide(21)
add_title_bar(s, 'v39-42：PV 分簇分层重建（动机）', '多 B 事件拆成单 B，单 B 处理路径完全不变')
add_bullets(s, [
    '背景：v38 对单 B 事件已经很好（无干扰），瓶颈在多 B/高连通事件——全图 91-139 条径迹，跨链边互相干扰',
    '核心想法（用户提出）：用 cluster 把多 B 事件按 PV 拆成多个单 B 子图，每个子图走与 v38 单 B 完全相同的处理路径（GNN/损失/重建都不变），只改怎么切',
    '关键要求：分簇器本身必须是一个带训练的 MLP（不能是真值硬切，否则训练/推理不一致）',
    '演进：',
    ('v39：先用 truth PV 硬切做训练侧试验 → EarlyStopping 提前结束（分簇器不可训练）', 1),
    ('v40：新增可训练 pv_cluster_head + 温度退火 → val 发散（train 切子图、val 全图）', 1),
    ('v41：课程学习 + val 对齐 + 子图 cap → 稳定但 3.2h/epoch', 1),
    ('v42：缩小规模验证 → 结论：子图训练伤模型（见后）', 1),
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S22 Gumbel ============
s = new_slide(22)
add_title_bar(s, 'v40：Gumbel-Softmax 温度退火（重参数化技巧）', '怎么让离散的分簇分配可训练')
add_bullets(s, [
    '问题：给每条径迹分配一个 PV（离散选择），argmax 不可导，梯度传不进去',
    'Gumbel 技巧（重参数化）：在 logit 上加 Gumbel 噪声再取 argmax，等价于按 softmax 概率采样',
    ('g = −log(−log(u))，u ~ Uniform(0,1)（Gumbel 分布采样）', 1),
    ('分配 = argmax( (logit + g) / τ )', 1),
    ('τ 大 → 近似按概率随机采样（探索，模型见到多样的分簇）；τ 小 → logit 主导，收敛到 argmax（确定）', 1),
    '温度退火：τ 从 1.0 → 0.1，与 B2 掩码退火同一个哲学——先探索后收敛',
    '配套：pv_cluster_head（可训练 MLP）吃 track/PV/边特征 → track 属 PV 的 logit，用与 pv_asso 相同的标签做 BCE 监督（独立训练）',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S23 课程学习 ============
s = new_slide(23)
add_title_bar(s, 'v41：课程学习（Curriculum Learning）+ 工程修复', '从易到难，避免随机分簇直接扰动好权重')
add_bullets(s, [
    '问题（v40 发散根因）：cluster 头随机初始化 + τ=0（resume 后绝对 epoch 使退火已结束）→ 一开始就是垃圾分簇，直接破坏 v38 的好权重',
    '课程学习：把训练顺序从难到易反过来',
    ('早期：alpha=1.0，全部用 truth PV 分簇（干净的单 B 子图，模型平滑适应子图结构）', 1),
    ('逐步过渡：alpha = 1 − run_epoch/30，每条径迹以概率 alpha 用 truth、否则用 cluster 头', 1),
    ('后期：alpha=0，全部用 cluster 头（与推理一致）', 1),
    '配套工程修复：',
    ('① val 也切子图（确定性分配）→ 消除 train/val 图结构 gap（v40 val 发散的主因）', 1),
    ('② 退火/课程相对本次 run 起点计（续训时绝对 epoch 已到 101，必须重置）', 1),
    ('③ 训练侧子图数量上限 cap（随机保留，提速）', 1),
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S24 v42 结论 ============
s = new_slide(24)
add_title_bar(s, 'v42：结论——子图训练伤模型', '为什么训练切子图这条路走不通')
add_table(s, [
    ['指标', 'v38（全图训练）', 'v42（子图训练 48ep）', '变化'],
    ['Perfect（thr0.9）', '29.26%', '26.28%', '-2.98pp'],
    ['LCAG class1 准确率', '76.8%', '56.38%', '-20.4pp'],
    ['val_combined', '35.86（全图口径）', '118.2（子图口径）', '口径不同，不可直接比'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(2.0), font_size=14)
add_bullets(s, [
    '根因：测试时 GNN 前向跑在全图上（重建流程：全图前向 → 剪枝 → 分簇 → 簇内重建），而训练让模型在子图上重学了 48 个 epoch → 全图前向能力退化（class1 掉 20pp）',
    '教训：train/infer 对齐必须看测试时模型到底在什么图上做前向，而不是只看 train/val 是否一致',
    '后续方向：退回全图训练（模型能力不退化）+ 推理时用分簇器拆多 B（分簇只是重建前的预处理，不改变 GNN 前向）',
], Inches(0.8), Inches(3.9), Inches(11.8), Inches(3.0), size=15)

# ============ S25 总结 ============
s = new_slide(25)
add_title_bar(s, '经验总结')
add_table(s, [
    ['#', '原则', '例子'],
    ['1', '先确认监督信号真的生效，再谈调模型', 'v31 class weight bug（拼写错误让权重根本没进损失）'],
    ['2', '离散决策放进训练 → 连续松弛 + 温度退火', 'B2 sigmoid 软掩码、Gumbel-Softmax 分簇'],
    ['3', '训练/推理对齐要看模型实际前向的图', '子图训练伤模型：测试前向仍跑全图'],
    ['4', '专项加权点到为止', 'class2 加权 3.0 → 2.0（3.0 伤 class1）'],
    ['5', '新头/新机制从 checkpoint 续训要重置状态', 'EarlyStopping、退火 epoch、课程 alpha 都要相对 run 起点'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(3.6), font_size=13)
add_bullets(s, [
    'PV 分簇路线最终结论（2026-08-24）：训练侧子图训练伤模型（class1 -20pp）；推理分簇也无效（Perfect 28.70% vs 29.26%，无显著差异）',
    '当前最优 = v38 全图方案（Perfect 29.26%）；后续优化另寻方向（如提升 class2、链合并、阈值扫描）',
], Inches(0.8), Inches(5.5), Inches(11.8), Inches(1.4), size=15)

# ============ S26 方向调整 ============
s = new_slide(26)
add_title_bar(s, '方向调整：评估口径转向 trigger（2026-08-25）', '用户明确：AllParticles 比 Perfect 更关键')
add_bullets(s, [
    '背景：导师汇报前，两条线并行（公开数据验证 v38 + PV 聚类基线）时，用户重新评估目标',
    '核心调整：',
    ('评估口径从 Perfect（完美重建整条链）转向 AllParticles（每个 B 的末态粒子尽量都重建到）——trigger 场景下"一个都不漏"比"全对"更重要', 1),
    ('区分 class1/class2 是物理难题（直系母子 vs 同母姊妹，需要精确的拓扑判别），优先级下调，不再作为紧迫优化目标', 1),
    '连带结论：凡是有损 AllParticles 的改动（如 PV 分簇切链），即使 Perfect 不掉也要重新评估是否值得',
    '待办：公开数据验证 v38 提升真实可靠 + 输出侧物理监督（见后续页）',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S27 输出侧物理监督动机 ============
s = new_slide(27)
add_title_bar(s, '输出侧物理监督：为什么选这个方向', '物理特征增强（改输入） vs 输出侧监督（加辅助头）')
add_bullets(s, [
    '先考虑过的方案——物理特征增强（改输入）：把手工算的 m_ij（两径迹不变质量）/cosθ 拼进边特征',
    ('问题：m_ij 对"是否同母"判别力弱（随机对也有类似分布），且引入人为物理假设，未必比 GNN 自己学的表示强', 1),
    ('用户质疑"这真的是一个好方法吗？" → 诚实分析后放弃', 1),
    '转向——输出侧监督（加辅助预测头）：保持端到端目标不变，用物理量作为辅助监督信号，逼主干表征携带物理信息',
    '文献支持（世界模型/JEPA 方向）：',
    ('HEP-JEPA (arXiv:2502.03933)、JetParticle-JEPA (arXiv:2606.14813)：预测任务 → 表示携带物理', 1),
    ('PhyIP (arXiv:2602.12218)：非侵入线性探针，验证"主干能否线性预测物理量"', 1),
    '落地方案：第7个监督头——边级 ππ 不变质量回归（mass head）',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S28 mass head 设计与首结果 ============
s = new_slide(28)
add_title_bar(s, 'mass head：边级不变质量回归（实现 + 首个结果）', '最小实验：从 v38 best 续训，验证"物理监督 → 更好重建"')
add_bullets(s, [
    '目标真值（无需额外标签，数据自带）：每条 tt 边两端轨迹的 ππ 不变质量',
    ('反归一化 px/py/pz（与 normalization_dict.pt 逐位核对）→ E = √(p² + m_pi²) → m² = (E₁+E₂)² − |p₁+p₂|²', 1),
    ('目标 = log10(m/MeV)（动态范围 2.4~4.3，避免原始质量 283-16959 MeV 压垮损失）', 1),
    ('掩码：未重建径迹（px≈py≈pz≈-1 哨兵）跳过；真实轨迹 px<0 占 48%，单维不能做掩码', 1),
    '头：输入 = decoder 边表征 latent_edges（16 维，op_trafo 覆盖前保存）→ MLP(16→64→1)；损失 SmoothL1(β=0.3)',
    '首个结果（仅 1 epoch，20 文件）：',
], Inches(0.8), Inches(1.3), Inches(11.8), Inches(3.2), size=14)
add_table(s, [
    ['指标', 'v38 基线', 'v38 + 1ep mass head', '变化'],
    ['AllParticles', '52.13%', '54.73%', '+2.6pp'],
    ['Perfect', '29.26%', '30.47%', '+1.2pp'],
    ['NoneIso', '47.87%', '45.27%', '-2.6pp'],
    ['mass_loss（train）', '—', '1.93 → 0.58（1 epoch）', '头在正常学习'],
], Inches(0.8), Inches(4.6), Inches(11.8), Inches(2.2), font_size=14)
add_bullets(s, [
    '注意：评估样本数有差异（12264 vs 12678），但 +2.6pp 方向性信号明确，完整 20 epoch 训练已重提',
], Inches(0.8), Inches(6.7), Inches(11.8), Inches(0.7), size=12)

# ============ S29 EarlyStopping bug ============
s = new_slide(29)
add_title_bar(s, 'EarlyStopping 续训 bug：只跑 1 个 epoch 就早停', 'checkpoint 里保存的早停状态在续训时被恢复')
add_bullets(s, [
    '现象：mass head 续训作业（9997828）第一个 epoch（val 38.95）后立即早停进入评估，只训练了 1 epoch',
    '根因：',
    ('v38 就是 EarlyStopping 结束的 → ckpt 里保存 wait_count=15（已满）、stopped_epoch=105、best_score=35.56', 1),
    ('续训恢复该状态后，第一个 epoch 的 val 未低于旧 best → wait_count 16 ≥ patience 15 → 立即触发早停', 1),
    '修复（exec_lightning.py）：新增 ResetEarlyStoppingOnResume 回调',
    ('resume 时在 fit 开始后把 wait_count 清零、best_score 重置为 inf → 早停从本 run 重新计', 1),
    '教训：与 v41 课程 alpha/退火 epoch 同类问题——所有"相对训练进度"的状态在续训时都要相对 run 起点重置',
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S30 PV 分簇 DA+重叠结论 ============
s = new_slide(30)
add_title_bar(s, 'PV 分簇（DA + 重叠软成员）评估结论', '确定性退火（DA）+ 允许重合的分簇 —— 对 trigger 目标不利，方向搁置')
add_table(s, [
    ['指标', 'v38 全图', 'DA 基线', 'DA + 重叠'],
    ['AllParticles', '52.13%', '49.12%', '46.41%'],
    ['Perfect', '29.26%', '28.27%', '28.20%'],
    ['NoneIso', '47.87%', '41.62%', '32.33%'],
    ['PartReco', '0%', '8.9%', '20.63%'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(2.6), font_size=14)
add_bullets(s, [
    '重叠软成员设计：score ≥ mem_thr（默认 0.3）的轨迹可同时属于多个 PV 簇，链只需在至少一个簇内完整；簇间用包含度判据去重',
    '结果解读：',
    ('重叠反而把链拆得更碎（PartReco 升到 20.6%），AllParticles 进一步下降 —— mem_thr=0.3 太松，轨迹被复制到多簇', 1),
    ('分簇切链这条路线整体与 trigger 目标（AllParticles 优先）相悖：两次评估均为负收益', 1),
    '结论：PV 分簇方向（DA/重叠/子图训练）全部搁置，重心转回全图方案 + 输出侧物理监督',
], Inches(0.8), Inches(4.5), Inches(11.8), Inches(2.6), size=15)

# ============ S31 公开数据验证 ============
s = new_slide(31)
add_title_bar(s, '公开数据验证 v38（进行中）', '用论文公开数据集（Run3 模拟, ~150 径迹/事件）训练 v38 算法栈，与论文结果对比')
add_bullets(s, [
    '目的：确认 v38 在 CERN 数据上的提升真实可靠，不是过拟合特定样本',
    '配置：50 trn 文件 / batch=1 / gacc=16（有效 16）/ lr=1e-3 / FP32 / epochs=100，v38 全部算法（B2、chain_lca、source_head、class2 加权）',
    '进度（version_45，2026-08-26 11:32）：',
    ('已跑 14 epoch（~1h/epoch），val_combined_loss 115.5 → 49.2（ep14 best，持续下降）', 1),
    ('val_LCA 0.92 → 0.60，val_tt 3.31 → 1.43 —— 训练健康', 1),
    ('⚠️ val_source_loss 从 ep1 起为 nan（公开数据 truth 链信息可能缺失），待查，不影响 combined', 1),
], Inches(0.8), Inches(1.4), Inches(11.8), Inches(5.6), size=15)

# ============ S32 当前状态 ============
s = new_slide(32)
add_title_bar(s, '当前状态与下一步', '2026-08-26')
add_table(s, [
    ['工作线', '状态', '说明'],
    ['公开数据验证 v38', '训练中（9995727, version_45）', '14/100 epoch，val 115→49，健康'],
    ['mass head 最小实验', '重提完整训练（10004201）', '1 epoch 已显示 AllParticles +2.6pp'],
    ['PV 分簇（DA/重叠）', '已搁置', '两次评估均伤 AllParticles'],
    ['PhyIP 线性探针', '待做', '验证主干表示能否线性预测物理量'],
], Inches(0.8), Inches(1.5), Inches(11.8), Inches(2.6), font_size=13)
add_bullets(s, [
    '下一步：',
    ('mass head 完整训练（20 epoch）跑完后评估，确认提升稳定', 1),
    ('训练后做 PhyIP 式线性探针：从主干节点/边表示线性回归 px/py/pz/m，对比 v38 与 v38+masshead 的可预测性', 1),
    ('公开数据 v38 训练完成后，与论文 DFEI 报告结果（arxiv:2304.08610）对比', 1),
], Inches(0.8), Inches(4.4), Inches(11.8), Inches(2.4), size=15)

prs.save(OUT)
print('[ok] 已保存: ' + OUT)
print('     共 ' + str(len(prs.slides._sldIdLst)) + ' 页')
