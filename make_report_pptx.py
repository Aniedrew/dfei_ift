#!/usr/bin/env python3
"""导师汇报 PPT: DFEI 近期成果 (直观版, 大字)
基于对话中讲透的直觉解释: B2 软剪刀 / 温度退火 / 源检测头 / chain_lca / PV分簇失败。
输出: report_figs/DFEI_progress_report.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = '/lzufs/home/guoqingxiang/dfei/scalable_mtl_hgnn'
FIG = BASE + '/report_figs'
OUT = FIG + '/DFEI_progress_report.pptx'
os.makedirs(FIG, exist_ok=True)

BLUE = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT = RGBColor(0xEA, 0xF2, 0xF8)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
BG_BOX = RGBColor(0xDD, 0xE8, 0xF3)
FONT = 'Microsoft YaHei'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.2), Inches(12.2), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = FONT
    ln = slide.shapes.add_shape(1, Inches(0.6), Inches(1.1), Inches(12.1), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()


def add_bullets(slide, items, left=0.8, top=1.4, width=11.8, height=5.8, size=21, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        txt, lvl = it if isinstance(it, tuple) else (it, 0)
        r = p.add_run()
        r.text = ('    ' * lvl) + ('•  ' if lvl == 0 else '–  ') + txt
        r.font.size = Pt(size - 2 * lvl)
        r.font.color.rgb = color or DARK
        r.font.name = FONT
    return tb


def add_table(slide, data, left, top, width, height, font_size=18, header_fill=BLUE):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for i in range(rows):
        for j in range(cols):
            c = tbl.cell(i, j)
            c.text = str(data[i][j])
            c.margin_left = Inches(0.05); c.margin_right = Inches(0.05)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.name = FONT
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i == 0 else DARK
                p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
            elif i % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = LIGHT
    return tbl


def add_box(slide, left, top, width, height, text, fill=BLUE, fg=RGBColor(0xFF, 0xFF, 0xFF), size=18, bold=True):
    sh = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for k, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = fg; r.font.name = FONT
    return sh


def add_arrow(slide, x1, y1, x2, y2):
    ln = slide.shapes.add_shape(8, Inches(x1), Inches(y1), Inches(max(x2 - x1, 0.05)), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = GRAY; ln.line.fill.background()


def footer(slide, idx):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.1), Inches(12.2), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = 'Qingxiang Guo · UCAS     ' + str(idx)
    r.font.size = Pt(12); r.font.color.rgb = GRAY; r.font.name = FONT


def new_slide(idx):
    s = prs.slides.add_slide(BLANK)
    footer(s, idx)
    return s


# ================= S1 封面 =================
s = new_slide(1)
add_box(s, 2.2, 1.6, 9.0, 1.6, 'DFEI 事件重建优化进展汇报', BLUE, size=36)
tb = s.shapes.add_textbox(Inches(2.2), Inches(3.5), Inches(9.0), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = '从 bug 修复到可微剪枝：Perfect 率 23.93% → 29.26%'
r.font.size = Pt(24); r.font.color.rgb = ORANGE; r.font.bold = True; r.font.name = FONT
tb = s.shapes.add_textbox(Inches(2.2), Inches(5.6), Inches(9.0), Inches(0.9))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = '汇报人：郭清祥 · 中国科学院大学 · 2026-08'
r.font.size = Pt(20); r.font.color.rgb = GRAY; r.font.name = FONT

# ================= S2 任务背景 =================
s = new_slide(2)
add_title(s, '任务：从噪声事件里找出 B 介子衰变链')
add_bullets(s, [
    '一次对撞产生几十到上百条径迹，绝大多数是背景',
    '目标：找出属于同一条 B 衰变链的径迹，重建衰变树',
    ('图：径迹=节点，径迹对=边；GNN 多任务学习（分类+剪枝+关联）', 1),
    ('难点：结构边只占 0.1%，极度稀有', 1),
    '核心指标：Perfect 率（链上所有粒子都重建正确的事件比例）',
], size=22)

# ================= S3 核心成果 =================
s = new_slide(3)
add_title(s, '核心成果：Perfect 率稳步提升')
add_box(s, 0.9, 1.6, 2.6, 2.0, 'v31 基线\n23.93%', RGBColor(0x8E, 0x8E, 0x8E), size=24)
add_box(s, 3.9, 1.6, 2.6, 2.0, 'v36 可微剪枝\n26.28%', BLUE, size=24)
add_box(s, 6.9, 1.6, 2.6, 2.0, 'v37 链判据\n27.25%', BLUE, size=24)
add_box(s, 9.9, 1.6, 2.6, 2.0, 'v38 当前最优\n29.26%', GREEN, size=24)
add_arrow(s, 3.5, 2.4, 3.9, 2.4)
add_arrow(s, 6.5, 2.4, 6.9, 2.4)
add_arrow(s, 9.5, 2.4, 9.9, 2.4)
add_bullets(s, [
    '总计提升 +5.33 个百分点（约 +22% 相对提升）',
    '每一轮都是小步改动，验证有效后再进下一轮',
], top=4.0, size=22)

# ================= S4 优化全景 =================
s = new_slide(4)
add_title(s, '优化路线全景：每一轮做了什么')
add_table(s, [
    ['版本', '改动', '一句话作用'],
    ['v31', '修 class weight bug', '逆向频率加权，让模型开始学结构边'],
    ['v36', 'B2 可微剪枝 + 温度退火', '训练时预演推理的剪枝'],
    ['v36', '源检测头', '让模型学会衰变树的形状'],
    ['v37', 'class2 加权', '专攻"同母边"这个结构瓶颈'],
    ['v38', '链内 LCA 一致性', '逼真链的边高置信且判对'],
], 0.8, 1.5, 11.8, 4.4, font_size=20)
add_bullets(s, [
    '还尝试了 PV 分簇（v40-42），验证后证明无效（后面详述）',
], top=6.2, size=20)

# ================= S5 B2 问题 =================
s = new_slide(5)
add_title(s, 'B2 可微剪枝：解决"训练图 ≠ 推理图"')
add_box(s, 0.9, 1.7, 5.2, 2.4, '训练时：全图\n每条边都参与消息传递\n（背景边也"说话"）', BLUE, size=20)
add_box(s, 7.2, 1.7, 5.2, 2.4, '推理时：先剪枝\n分数低于阈值 0.9 的边\n被彻底删除', ORANGE, size=20)
add_arrow(s, 6.1, 2.7, 7.2, 2.7)
tb = s.shapes.add_textbox(Inches(5.4), Inches(2.35), Inches(2.2), Inches(0.8))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = '结构不一致'
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = RED; r.font.name = FONT
add_bullets(s, [
    '模型训练时习惯了"每条边都掺一点信息"',
    '推理时低分边被删了，模型"水土不服"',
    'B2 的思路：训练时就把推理的剪枝"预演"一遍',
], top=4.6, size=22)

# ================= S6 B2 解法 =================
s = new_slide(6)
add_title(s, 'B2 软剪刀：让低分边"不说话"')
add_bullets(s, [
    '消息传递本来按边分数加权：低分边贡献本来就小',
    'B2 加一道"软剪刀"：分数低于阈值(0.85)的边，',
    ('把它的发言权再压到接近 0 —— 就像推理时把它删掉', 1),
    '但"软"的：还留一点梯度能传回来，模型才能学',
    '只在最后一个 GN block 启用（与推理剪枝位置对齐）',
], size=22)
add_box(s, 0.9, 5.2, 11.6, 1.4, '一句话：训练时把"该剪的边"按推理的方式处理，让模型在稀疏图上学习', BG_BOX, fg=BLUE, size=22, bold=False)

# ================= S7 温度退火 =================
s = new_slide(7)
add_title(s, '温度退火：从"钝剪刀"到"利剪刀"')
add_box(s, 0.9, 1.7, 5.4, 2.2, '训练初期：τ 大（钝）\n低分边还留一点余量\n模型先学大方向', GREEN, size=20)
add_box(s, 7.0, 1.7, 5.4, 2.2, '训练后期：τ 小（利）\n几乎等于硬剪枝\n逼近真实推理', RED, size=20)
add_arrow(s, 6.3, 2.6, 7.0, 2.6)
add_bullets(s, [
    '如果一开始就"一刀切"：被剪的边一点梯度都没有，模型学不动',
    '所以先模糊后精确——和模拟退火一个思路',
    '类比：先在有点吵的教室讲课，考试时是安静的教室',
], top=4.4, size=22)

# ================= S8 源检测头 =================
s = new_slide(8)
add_title(s, '源检测头：让模型学会"树的形状"')
add_bullets(s, [
    'Rumor Centrality（RC）= 老师：对每条真链算"拓扑根"',
    ('即"链最顶端的节点"，只靠图结构，不用物理量', 1),
    'source_head = 学生：从节点特征猜"谁是根"',
    ('这是个附加题，推理时并不考它', 1),
    '但为了答对，主干必须学会"衰变链谁在谁上面"',
    ('这种层级理解，正好是 LCA 分类和链重建需要的', 1),
], size=22)
add_box(s, 0.9, 5.6, 11.6, 1.3, 'RC 提供"根在哪"的答案，模型通过"猜根"被迫学会树的结构', BG_BOX, fg=BLUE, size=21, bold=False)

# ================= S9 class2 + chain_lca =================
s = new_slide(9)
add_title(s, '专攻瓶颈：class2 与链内一致性')
add_bullets(s, [
    'class2（同母边）是最大结构瓶颈：被判对只有 ~45%',
    ('一半被判成 class1（母子），混淆严重', 1),
    '招数一：class2 专项加权（放大它的损失，点到为止）',
    '招数二：链内 LCA 一致性（侦探拼拼图）：',
    ('真链的边必须"高置信"（hinge 惩罚低置信边）', 1),
    ('还必须"判对类别"（额外 CE 监督）', 1),
], size=21)

# ================= S10 PV 分簇 =================
s = new_slide(10)
add_title(s, '一次验证失败的尝试：PV 分簇（v40-42）')
add_bullets(s, [
    '动机：把多 B 事件按 PV 拆成单 B 子图，保留单 B 的优点',
    ('分簇器做成可训练 MLP + Gumbel + 课程学习，做了很多工程', 1),
    '结果 1（训练侧切子图）：模型反而变差',
    ('LCAG class1 准确率 76.8% → 56.4%，Perfect -3pp', 1),
    ('根因：测试时模型前向仍跑在全图上，子图训练白费', 1),
    '结果 2（推理侧分簇）：Perfect 28.70% vs 29.26%，无提升',
    '结论：路线关闭，但教训宝贵（见下页）',
], size=20)

# ================= S11 经验总结 =================
s = new_slide(11)
add_title(s, '经验总结：五条原则')
add_table(s, [
    ['#', '原则', '例子'],
    ['1', '先确认监督信号生效，再调模型', 'class weight 拼写 bug'],
    ['2', '离散决策进训练 → 连续松弛+退火', 'B2 软剪刀、Gumbel'],
    ['3', '对齐要看"模型实际前向的图"', '子图训练伤模型'],
    ['4', '专项加权点到为止', 'class2 加权 3.0→2.0'],
    ['5', '续训要重置调度状态', 'EarlyStopping/退火'],
], 0.8, 1.5, 11.8, 4.4, font_size=20)
add_bullets(s, [
    '最贵的教训：训练/推理对齐，不是"train/val 一致"就行',
    ('要看测试时模型到底在什么图上做前向', 1),
], top=6.0, size=20)

# ================= S12 当前与下一步 =================
s = new_slide(12)
add_title(s, '当前状态与下一步')
add_box(s, 0.9, 1.6, 5.4, 1.6, '当前最优：v38 全图方案\nPerfect 29.26%', GREEN, size=22)
add_box(s, 7.0, 1.6, 5.4, 1.6, '下一步：专攻 class2\n（仍是链结构最大短板）', BLUE, size=22)
add_bullets(s, [
    '候选方案（按成本排序）：',
    ('① 推理侧衰变树一致性约束（零训练成本，先试）', 1),
    ('② 训练侧 class0 负采样 + focal loss（解决稀有性）', 1),
    ('③ 物理特征增强：同母对的不变质量（天花板最高）', 1),
], top=3.7, size=21)
add_box(s, 0.9, 5.9, 11.6, 1.2, '汇报完，感谢！', BG_BOX, fg=BLUE, size=22, bold=False)

prs.save(OUT)
print('[ok] 已保存: ' + OUT)
print('     共 ' + str(len(prs.slides._sldIdLst)) + ' 页')
