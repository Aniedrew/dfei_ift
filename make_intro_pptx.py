#!/usr/bin/env python3
"""DFEI 工作介绍 PPT（给不了解项目的老师）。

前半部分: DFEI 原版方法介绍（基于 scalable_mtl_hgnn 原版 + 论文 arXiv:2504.21844 / 2304.08610）
后半部分: 我的优化工作（v31 基线 -> v38 + v39 PV 分簇 + 推理侧探索 + 3D 可视化）
"""
import os
import io
import zstandard as zstd
import numpy as np
import torch
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from mpl_toolkits.mplot3d import Axes3D
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

BASE = "/lzufs/home/guoqingxiang/dfei/scalable_mtl_hgnn"
FIG = f"{BASE}/intro_figs"
os.makedirs(FIG, exist_ok=True)
OUT = f"{FIG}/DFEI_work_intro.pptx"

plt.rcParams["font.sans-serif"] = ["Noto Sans CJK SC", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False

BLUE = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xC0, 0x39, 0x2B)

# ================= 数据 =================
VERSIONS = ["v31 基线", "v36", "v37", "v38"]
PERFECT = [23.93, 26.28, 27.25, 29.26]
ALL = [43.42, 49.24, 50.60, 52.13]
NONEISO = [56.58, 50.76, 49.40, 47.87]

LCA_V = ["class0", "class1", "class2", "class3"]
LCA31 = [98.40, 67.78, 41.26, 75.43]
LCA38 = [97.86, 76.81, 47.92, 55.63]


# ================= 图片 1: 重建指标对比 =================
def fig_metrics():
    x = np.arange(len(VERSIONS))
    w = 0.26
    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    b1 = ax.bar(x - w, PERFECT, w, label="PerfectReco", color="#1f6fd6")
    b2 = ax.bar(x, ALL, w, label="AllParticles", color="#2ca02c")
    b3 = ax.bar(x + w, NONEISO, w, label="NoneIso (越低越好)", color="#d62728")
    for bars in (b1, b2, b3):
        for r in bars:
            ax.text(r.get_x() + r.get_width() / 2, r.get_height() + 0.6,
                    f"{r.get_height():.1f}", ha="center", fontsize=9)
    ax.set_xticks(x); ax.set_xticklabels(VERSIONS, fontsize=13)
    ax.set_ylabel("比例 (%)", fontsize=13)
    ax.set_ylim(0, 65)
    ax.legend(fontsize=12, loc="upper left")
    ax.grid(axis="y", alpha=0.3)
    ax.set_title("重建性能逐版本提升（thr 0.9 同口径, CERN 官方 MC）", fontsize=14)
    plt.tight_layout(); plt.savefig(f"{FIG}/fig_metrics.png", dpi=150); plt.close()


# ================= 图片 2: LCAG 分类对比 =================
def fig_lca():
    x = np.arange(len(LCA_V))
    w = 0.32
    fig, ax = plt.subplots(figsize=(8.5, 4.5))
    b1 = ax.bar(x - w / 2, LCA31, w, label="v31 基线", color="#8faadc")
    b2 = ax.bar(x + w / 2, LCA38, w, label="v38", color="#1f6fd6")
    for bars in (b1, b2):
        for r in bars:
            ax.text(r.get_x() + r.get_width() / 2, r.get_height() + 1.0,
                    f"{r.get_height():.1f}", ha="center", fontsize=9)
    ax.set_xticks(x); ax.set_xticklabels(["class0\n(无关边)", "class1\n(跨B边)",
                                          "class2\n(同B边)", "class3"], fontsize=12)
    ax.set_ylabel("准确率 (%)", fontsize=13)
    ax.set_ylim(0, 110)
    ax.legend(fontsize=12)
    ax.grid(axis="y", alpha=0.3)
    ax.set_title("LCAG 层级分类准确率（结构瓶颈的改善）", fontsize=14)
    plt.tight_layout(); plt.savefig(f"{FIG}/fig_lca.png", dpi=150); plt.close()


# ================= 图片 3: 重建流程示意图 =================
def fig_pipeline():
    fig, ax = plt.subplots(figsize=(11.5, 3.0))
    ax.axis("off")
    boxes = [
        (0.02, "事件图\ntracks/pvs 节点\n+ 关联边", "#eef4fb", "#1f4e79"),
        (0.24, "HGNN 多任务\nLCA 分类 · 节点/边剪枝 · PV 关联", "#eef7ee", "#2ca02c"),
        (0.46, "阈值剪枝\nnode/edge 过滤", "#fdf6ec", "#d98100"),
        (0.68, "LCAG 树聚类\n衰变链重建", "#fdeef0", "#d62728"),
        (0.90, "链级评估\nPerfect / All / NoneIso", "#f4f0fb", "#7b2cbf"),
    ]
    for fx, txt, fc, ec in boxes:
        ax.add_patch(FancyBboxPatch((fx, 0.18), 0.15, 0.62, boxstyle="round,pad=0.02",
                                    fc=fc, ec=ec, lw=1.6))
        ax.text(fx + 0.075, 0.49, txt, ha="center", va="center", fontsize=10.5, color=ec)
        if fx < 0.9:
            ax.add_patch(FancyArrowPatch((fx + 0.155, 0.49), (fx + 0.235, 0.49),
                                         arrowstyle="-|>", mutation_scale=18, color="#888"))
    ax.set_xlim(0, 1.08); ax.set_ylim(0, 1)
    plt.tight_layout(); plt.savefig(f"{FIG}/fig_pipeline.png", dpi=150); plt.close()


# ================= 图片 4: 事件 3D 示意图 =================
def fig_event3d():
    DATA_FILE = "/lzufs/user/guoqingxiang/DFEI_IFT_20260702/data/MC_normed/inclusive_00342442/tst_data_00231000_00231999.pt.zst"
    NORM = "/lzufs/user/guoqingxiang/DFEI_IFT_20260702/dfei_repo/preprocessing/normalization_dict.pt"
    dctx = zstd.ZstdDecompressor()
    with open(DATA_FILE, "rb") as f:
        with dctx.stream_reader(f) as r:
            data = torch.load(io.BytesIO(r.read()), weights_only=False)
    evt = data[747]
    x = evt["tracks"]["x"].numpy()
    norm = torch.load(NORM, map_location="cpu", weights_only=False)
    c, s = norm["center"], norm["scale"]
    names = ["px_reco", "py_reco", "pz_reco", "xProd_reco", "yProd_reco", "zProd_reco"]
    cols = {nm: i for i, nm in enumerate(names)}
    phys = np.zeros((x.shape[0], 6))
    for nm in names:
        phys[:, cols[nm]] = x[:, cols[nm]] * s[nm] + c[nm]
    px, py, pz = phys[:, 0], phys[:, 1], phys[:, 2]
    x0, y0, z0 = phys[:, 3] / 1000, phys[:, 4] / 1000, phys[:, 5] / 1000
    CHAIN0 = {1, 13, 14, 24, 51, 52, 69}
    CHAIN1 = {7, 19, 46, 58}

    fig = plt.figure(figsize=(8.5, 6.2))
    ax = fig.add_subplot(111, projection="3d")
    L = 0.10
    for i in range(x.shape[0]):
        d = np.array([px[i], py[i], pz[i]], dtype=float)
        dn = d / (np.linalg.norm(d) + 1e-9)
        end = np.array([x0[i], y0[i], z0[i]]) + dn * L
        if i in CHAIN0:
            col, lw, alpha = "#e63946", 3.5, 1.0
        elif i in CHAIN1:
            col, lw, alpha = "#1d6fd6", 3.5, 1.0
        else:
            col, lw, alpha = "#9aa0a6", 1.2, 0.30
        ax.plot([x0[i], end[0]], [y0[i], end[1]], [z0[i], end[2]],
                color=col, lw=lw, alpha=alpha)
        if i in CHAIN0 or i in CHAIN1:
            ax.scatter([x0[i]], [y0[i]], [z0[i]], color=col, s=25, alpha=1.0)
    ax.set_xlim(-0.05, 0.05); ax.set_ylim(-0.05, 0.05); ax.set_zlim(-0.10, 0.45)
    ax.set_xlabel("x (m)"); ax.set_ylabel("y (m)"); ax.set_zlabel("z (束流, m)")
    ax.set_title("事件 92441664 · 双B衰变 · 红=链0 (v31失败/v37成功), 蓝=链1, 灰=背景", fontsize=12)
    plt.tight_layout(); plt.savefig(f"{FIG}/fig_event3d.png", dpi=150); plt.close()


# ================= 生成图片 =================
fig_metrics()
fig_lca()
fig_pipeline()
fig_event3d()
print("[ok] 图片已生成")

# ================= PPT =================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def title_bar(s, text):
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = BLUE
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(1.02), Inches(12.1), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()


def bullets(s, items, left=Inches(0.8), top=Inches(1.4), width=Inches(11.7),
            height=Inches(5.0), size=18, color=DARK):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        if isinstance(it, tuple):
            r = p.add_run(); r.text = "\u2022  " + it[0]
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = BLUE
            p2 = p  # 保持简单: 同段追加说明
            r2 = p.add_run(); r2.text = " — " + it[1]
            r2.font.size = Pt(size - 2); r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = "\u2022  " + it
            r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def pic(s, fname, left, top, width=None, height=None):
    from PIL import Image
    path = os.path.join(FIG, fname)
    if not os.path.exists(path):
        return
    with Image.open(path) as im:
        w, h = im.size
    ar = w / h
    if width is None and height is None:
        width = Inches(9.0)
    if width is not None and height is None:
        height = width / ar
    s.shapes.add_picture(path, left, top, width=width, height=height)


def table(s, data, left, top, width, height, col_widths=None, font=14, header_fill="1F4E79"):
    from pptx.util import Inches as IN
    rows, cols = len(data), len(data[0])
    gt = s.shapes.add_table(rows, cols, left, top, width, height).table
    for j in range(cols):
        if col_widths:
            gt.columns[j].width = col_widths[j]
    for i in range(rows):
        for j in range(cols):
            cell = gt.cell(i, j)
            cell.text = str(data[i][j])
            cell.text_frame.paragraphs[0].font.size = Pt(font)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i == 0 else DARK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(header_fill)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    return gt


# ===== Slide 1 封面 =====
s = new_slide()
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.3))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "基于异构图神经网络的 LHCb 全事件重建"
r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = BLUE
tb = s.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.8))
tf = tb.text_frame
r = tf.paragraphs[0].add_run(); r.text = "DFEI/HGNN：原版方法介绍与我的优化工作"
r.font.size = Pt(28); r.font.color.rgb = DARK
tb = s.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.3), Inches(0.6))
tf = tb.text_frame
r = tf.paragraphs[0].add_run(); r.text = "Qingxiang Guo · University of Chinese Academy of Sciences"
r.font.size = Pt(18); r.font.color.rgb = GRAY
tb = s.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.6))
tf = tb.text_frame
r = tf.paragraphs[0].add_run(); r.text = "数据: CERN 官方 MC · GPU: RTX 2080 Ti (10 GB)"
r.font.size = Pt(16); r.font.color.rgb = GRAY

# ===== 前半部分: DFEI =====
# Slide 2 背景
s = new_slide(); title_bar(s, "背景：高亮度下的 LHCb 重建挑战")
bullets(s, [
    ("LHCb 实验", "专攻 b/c 强子衰变的精密测量（CP 破坏、稀有衰变、味物理）"),
    ("Upgrade I (Run 3/4)", "平均 ~5 次 pp 碰撞/事件、~150 条径迹；Upgrade II 将到 ~50 碰撞、~1000 径迹"),
    ("三重挑战", "触发低延迟（O(100 ms)/事件）、存储受限（O(10 PB)/年）、背景组合爆炸"),
    ("传统 exclusive 选择", "逐衰变道找候选，忽略事件上下文信息，难以应对高复杂度"),
    ("PV 误关联", "高亮度下径迹归属错误 → 影响飞行距离与 CP 测量精度"),
], size=20)

# Slide 3 DFEI 思想
s = new_slide(); title_bar(s, "DFEI：全事件解释")
bullets(s, [
    ("核心思想", "用一个模型同时识别、隔离、层级重建每事件中的所有重味强子衰变链"),
    ("GNN 适配性", "图结构天然表示事件：变长、稀疏、异构，无需人工排序"),
    ("与 FEI 对比", "Belle II 的 FEI 为每个衰变道训练分类器；DFEI 单模型统括，并引入上下文"),
    ("LCAG 层级目标", "Kahn 等提出：用\"最低共同祖先代数\"标签学习衰变层级结构"),
    ("原版 DFEI (多阶段)", "节点剪枝 GNN → 边剪枝 GNN → LCAG 分类 GNN，逐级过滤背景"),
], size=20)

# Slide 4 事件图表示
s = new_slide(); title_bar(s, "事件 → 异构图表示")
bullets(s, [
    ("节点", "tracks（径迹）与 pvs（主顶点）两类节点"),
    ("边", "tracks-tracks（径迹关联）、tracks-pvs（顶点关联）"),
    ("特征", "动量 (px,py,pz)、产生点 (xProd,yProd,zProd)、电荷、PID 概率、ghost 概率等"),
    ("规模", "~150 径迹的事件 → 上万条候选边"),
], top=Inches(1.4), size=20)
pic(s, "fig_pipeline.png", Inches(1.0), Inches(4.2), width=Inches(11.0))

# Slide 5 HGNN 原版方法
s = new_slide(); title_bar(s, "HGNN 多任务架构（原版 scalable_mtl_hgnn）")
bullets(s, [
    ("编码器", "各节点/边类型独立 MLP 编码（保留异构信息）"),
    ("GN 块", "多轮图消息传递，集成可微图剪枝 + 加权消息传递"),
    ("解码器/多任务头", "LCA 分类（4 类）、节点剪枝、边剪枝、PV 关联 —— 联合训练共享表征"),
    ("训练", "CERN 模拟环境（DFEI_IFT 数据），多任务损失加权"),
    ("对应论文", "Sutcliffe et al., arXiv:2504.21844；DFEI 原型 García Pardiñas et al., arXiv:2304.08610"),
], size=20)

# Slide 6 重建后处理与评估指标
s = new_slide(); title_bar(s, "推理流程与评估指标")
bullets(s, [
    ("阈值剪枝", "按 node/edge 置信度阈值过滤背景节点与边"),
    ("LCAG 树聚类", "由 LCA 层级标签聚类出衰变链（reconstruct_decay）"),
    ("PerfectReco", "链完整且无背景混入（最高要求）"),
    ("AllParticles", "所有信号粒子都在（可有污染）"),
    ("NoneIso", "无污染（可有缺失）"),
], size=20)
pic(s, "fig_metrics.png", Inches(1.5), Inches(3.6), width=Inches(10.0))

# ===== 后半部分: 我的工作 =====
# Slide 7 概览
s = new_slide(); title_bar(s, "我的工作概览")
bullets(s, [
    ("① 基线建立", "复现原版 → 定位并修复 class-weight bug → CERN 基线 v31"),
    ("② 优化 A", "B2 可微剪枝：消除训练-推理剪枝鸿沟（v36）"),
    ("③ 优化 B", "源检测头：Rumor Centrality 训练化，学\"根-叶\"结构"),
    ("④ 优化 C", "loss 再平衡：class2 加权 + 链级 LCA 一致性 + 链内类别 CE（v37/v38）"),
    ("⑤ 优化 D", "PV 分簇分层重建（v39，训练中）"),
    ("⑥ 推理侧探索", "chain_lca_filter 后处理判据验证（结论：区分度不足，已放弃）"),
], size=20)

# Slide 8 基线 bug
s = new_slide(); title_bar(s, "① 基线建立：class-weight bug 修复")
bullets(s, [
    ("问题", "配置键拼写错误（LCA__weights → LCA_weights），类别权重未生效"),
    ("后果", "LCAG class1 准确率跌至 ~0%，PerfectReco 严重受限"),
    ("修复后", "CERN 官方 MC 上建立基线 v31：Perfect 23.9%、NoneIso 56.6%"),
], size=20)
pic(s, "fig_lca.png", Inches(1.5), Inches(3.7), width=Inches(10.0))

# Slide 9 优化 A
s = new_slide(); title_bar(s, "② 优化 A：B2 可微剪枝（消除 train-inference gap）")
bullets(s, [
    ("痛点", "训练用全图监督，推理时硬阈值剪枝 → 分布割裂（gap）"),
    ("方案", "训练时对 node/edge 权重施加软掩码 σ((w−cut)/τ)，温度退火 1.0→0.1"),
    ("效果", "训练即模拟剪枝，收敛到与推理一致的稀疏结构"),
    ("结果 (v36)", "Perfect 26.28%（+2.35pp vs v31）、NoneIso 50.8%（−5.8pp）"),
], size=20)

# Slide 10 优化 B
s = new_slide(); title_bar(s, "③ 优化 B：源检测头（Rumor Centrality 训练化）")
bullets(s, [
    ("物理直觉", "真实衰变链是\"有根树\"（B 介子为根），噪声团块无根"),
    ("推理侧", "用 rumor centrality 找链根（B 介子候选）"),
    ("训练侧", "新增第 6 个监督头：预测每条 truth 链的根节点（BCE）"),
    ("作用", "让主干显式学习\"根-叶\"层级结构，与推理侧 RC 对齐"),
], size=20)

# Slide 11 优化 C
s = new_slide(); title_bar(s, "④ 优化 C：Loss 再平衡（结构瓶颈专项）")
bullets(s, [
    ("方案4 class2 加权", "同B边（class2）是链内结构瓶颈（~41%），专项放大其 CE 权重"),
    ("方案5 链级 LCA 一致性", "鼓励真链内边\"高置信\"（chain_lca_filter 训练化）"),
    ("方案6 b2_cut 对齐 + 链内类别 CE", "剪枝阈值 0.5→0.85 对齐推理；链内边额外 CE 监督类别正确"),
    ("结果 (v38)", "Perfect 29.26%（vs v37 +2.0pp）；class1 76.8%、class2 47.9%"),
], size=20)
pic(s, "fig_lca.png", Inches(1.5), Inches(4.2), width=Inches(10.0))

# Slide 12 优化 D
s = new_slide(); title_bar(s, "⑤ 优化 D：PV 分簇分层重建（v39，训练中）")
bullets(s, [
    ("问题", "CERN 事件 91–139 tracks，高连通图 → 跨链干扰（2B/>2B 事件是重灾区）"),
    ("方案", "按 PV 把事件分成簇（每簇 20–30 节点），簇内独立重建衰变链再合并"),
    ("实现", "reconstruction.py 新增 pv_cluster 逻辑；真/预测 PV 分配可选"),
    ("现状", "v39 训练已启动（resume v38 best），预计 14h 完成"),
], size=20)

# Slide 13 推理侧探索
s = new_slide(); title_bar(s, "⑥ 推理侧探索：chain_lca_filter 后处理")
bullets(s, [
    ("想法", "用链级 LCA 置信度阈值过滤\"物理不自洽\"的链，压低 NoneIso"),
    ("验证", "v36/v37 record 数据上扫阈值（conf / class2_frac / 联合）"),
    ("结论", "Perfect 链 conf 中位 0.65 vs 失败链 0.55，重叠过大——过滤误杀真链，任何阈值都掉点"),
    ("教训", "链级平均置信度判据无区分度；需更细粒度判据或 selection MLP 链打分"),
], size=20)

# Slide 14 结果总览
s = new_slide(); title_bar(s, "结果总览：v31 → v38（thr 0.9 同口径）")
table(s, [
    ["版本", "PerfectReco", "AllParticles", "NoneIso", "LCAG class1", "LCAG class2"],
    ["v31 基线", "23.93%", "43.42%", "56.58%", "67.78%", "41.26%"],
    ["v36", "26.28%", "49.24%", "50.76%", "64.51%", "47.66%"],
    ["v37", "27.25%", "50.60%", "49.40%", "69.14%", "44.22%"],
    ["v38", "29.26%", "52.13%", "47.87%", "76.81%", "47.92%"],
], Inches(0.9), Inches(1.5), Inches(11.5), Inches(2.6), font=15)
pic(s, "fig_metrics.png", Inches(1.7), Inches(4.3), width=Inches(10.0))

# Slide 15 3D 可视化
s = new_slide(); title_bar(s, "3D 可视化：衰变级联视图")
pic(s, "fig_event3d.png", Inches(2.0), Inches(1.4), width=Inches(9.0))
tb = s.shapes.add_textbox(Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "事件 92441664（双B、7 粒子长链）：v31 重建失败 → v37 重建成功；含 PV、衰变级联边、粒子类型"
r.font.size = Pt(14); r.font.color.rgb = GRAY

# Slide 16 总结与展望
s = new_slide(); title_bar(s, "总结与展望")
bullets(s, [
    ("成果", "v31→v38：PerfectReco +5.3pp（23.9%→29.3%），NoneIso 首次跌破 48%"),
    ("方法沉淀", "B2 可微剪枝 + 源检测头 + 链级结构 loss 的组合被验证有效"),
    ("进行中", "v39 PV 分簇分层重建（降单图复杂度，目标 2B/>2B 事件）"),
    ("展望", "selection MLP 候选链打分（需训练 scorer）、更细粒度链级判据"),
    ("其他产出", "GPU 调度/效率分析笔记、版本演进记录、3D 可视化工具"),
], size=20)

prs.save(OUT)
print(f"[ok] PPT 已生成: {OUT}")
