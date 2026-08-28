"""生成简单的 DFEI 组会 PPT (基于 meeting_20260811_slides.md)
输出到 meeting_20260811_figs/DFEI_progress_20260811.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

BASE = "/lzufs/home/guoqingxiang/dfei/scalable_mtl_hgnn"
FIG = f"{BASE}/meeting_20260811_figs"
OUT = f"{FIG}/DFEI_progress_20260811.pptx"

BLUE = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x2C, 0xA0, 0x2C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def add_title_bar(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = BLUE
    # 下划线
    ln = slide.shapes.add_shape(1, Inches(0.6), Inches(1.02), Inches(12.1), Pt(3))
    ln.fill.solid()
    ln.fill.fore_color.rgb = BLUE
    ln.line.fill.background()

def add_bullets(slide, items, left, top, width, height, size=16):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = ("\u2022  " if isinstance(it, str) else "\u2022  ") + it
        r.font.size = Pt(size)
        r.font.color.rgb = DARK
    return tb

def add_pic(slide, fname, left, top, width=None, height=None):
    path = os.path.join(FIG, fname)
    if not os.path.exists(path):
        return
    with Image.open(path) as im:
        w, h = im.size
    ar = w / h
    if width is None and height is None:
        width = Inches(7.2)
    if width is not None and height is None:
        height = width / ar
    if height is not None and width is None:
        width = height * ar
    slide.shapes.add_picture(path, left, top, width=width, height=height)

def add_footer(slide, idx):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Qingxiang Guo \u00b7 UCAS \u00b7 DFEI group meeting \u00b7 Aug 11, 2026      {idx}"
    r.font.size = Pt(11)
    r.font.color.rgb = GRAY

def new_slide(idx):
    s = prs.slides.add_slide(BLANK)
    add_footer(s, idx)
    return s

# ============ Slide 1: Title ============
s = new_slide(1)
tb = s.shapes.add_textbox(Inches(1.2), Inches(2.4), Inches(10.9), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "DFEI Local Run Results and Optimization Strategy"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = BLUE
tb = s.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(10.9), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Progress Report"
r.font.size = Pt(30); r.font.color.rgb = DARK
tb = s.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10.9), Inches(0.7))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "official CERN MC & public dataset \u00b7 RTX 2080 Ti (10 GB)"
r.font.size = Pt(18); r.font.color.rgb = GRAY
tb = s.shapes.add_textbox(Inches(1.2), Inches(5.4), Inches(10.9), Inches(0.6))
tf = tb.text_frame
r = tf.paragraphs[0].add_run()
r.text = "Qingxiang Guo \u00b7 University of Chinese Academy of Sciences \u00b7 DFEI group meeting \u00b7 Aug 11, 2026"
r.font.size = Pt(16); r.font.color.rgb = GRAY

# ============ Slide 2: What This Report Covers ============
s = new_slide(2)
add_title_bar(s, "What This Report Covers")
add_bullets(s, [
    "Reconstruction results on the official CERN MC and the public dataset",
    "A concrete training-speed problem on the 10 GB GPU",
    "Why hard-threshold pruning loses signal chains (diagnosis)",
    "Three targeted changes to fix it (edges / nodes / candidate selection)",
], Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5), size=20)

# ============ Slide 3: Setup ============
s = new_slide(3)
add_title_bar(s, "Setup")
rows = [
    ("Code", "scalable_mtl_hgnn (branch yukai_IFT, weight-fix carried locally)"),
    ("Data 1", "Official CERN MC \u2014 DFEI_IFT_20260702/MC_normed (~91 trk/evt)"),
    ("Data 2", "Public (paper) dataset \u2014 converted_LHCbcollision (~139 trk/evt)"),
    ("GPU", "NVIDIA RTX 2080 Ti, 10.6 GB VRAM"),
    ("Model", "HGNN, 4 message-passing blocks, ~570 k params"),
]
left, top, w, h = Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.9)
for i, (k, v) in enumerate(rows):
    y = top + Inches(0.95 * i)
    tb = s.shapes.add_textbox(left, y, Inches(1.8), h)
    tf = tb.text_frame
    r = tf.paragraphs[0].add_run(); r.text = k
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = BLUE
    tb = s.shapes.add_textbox(left + Inches(2.1), y, Inches(9.4), h)
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = v
    r.font.size = Pt(18); r.font.color.rgb = DARK

# ============ Slide 4: Reconstruction Results ============
s = new_slide(4)
add_title_bar(s, "Reconstruction Results (thr 0.2)")
add_bullets(s, [
    "CERN MC (v25): per-chain Perfect 27.7%, AllParticles 55.6%, NoneIso 6.0%, PartReco 15.9%, NotFound 22.5%",
    "Public (v23): per-chain Perfect 19.8%",
    "Paper WHGNN reference: per-event Perfect 21.5% (different metric def., harder test set)",
    "Status: retraining with corrected class weights in progress",
], Inches(0.8), Inches(1.35), Inches(6.9), Inches(5.0), size=15)
add_pic(s, "fig07_reco_categories.png", Inches(7.9), Inches(1.35), width=Inches(5.1))

# ============ Slide 5: Problem 1 - Training Is Slow ============
s = new_slide(5)
add_title_bar(s, "Problem 1: Training Is Slow")
add_bullets(s, [
    "~30\u201350 min/epoch on 10 GB RTX 2080 Ti (full graph: ~150 tracks \u2192 ~23k edges/event)",
    "10 GB VRAM forces small batches \u2014 the main slowdown",
    "100 epochs \u2248 days",
], Inches(0.8), Inches(1.35), Inches(6.9), Inches(3.0), size=16)
add_pic(s, "fig10_epoch_time.png", Inches(7.9), Inches(1.35), width=Inches(5.1))
add_pic(s, "fig02_training_loss.png", Inches(7.9), Inches(4.4), width=Inches(5.1))

# ============ Slide 6: Problem 2 - Pruning Loses Chains ============
s = new_slide(6)
add_title_bar(s, "Problem 2: Hard-Threshold Pruning Loses Chains")
add_bullets(s, [
    "Perfect \u2248 (chain survives pruning) \u00d7 (structure correct)",
    "At thr 0.2, pruning kills ~30% of chains via edges and ~26% via nodes (overlap; edges dominate)",
    "Only ~40% of surviving chains get the exact LCA structure right",
    "Two bottlenecks: greedy hard threshold, and classification quality",
], Inches(0.8), Inches(1.35), Inches(6.9), Inches(5.0), size=16)
add_pic(s, "fig06_chain_survival.png", Inches(7.9), Inches(1.35), width=Inches(5.1))

# ============ Slide 7: Pruning Heads Are Strong ============
s = new_slide(7)
add_title_bar(s, "Pruning Heads Are Strong \u2014 It's the Threshold, Not the Heads")
add_bullets(s, [
    "Edge-pruning AUC: 0.999 (CERN) / 1.000 (public)",
    "Node-pruning AUC: 0.974 / 0.997",
    "Ranking is nearly perfect \u2014 yet the 0.2 cut still deletes true edges",
    "Data note: public data has LCA class-2/3 edges; CERN MC almost none",
], Inches(0.8), Inches(1.35), Inches(6.9), Inches(4.0), size=15)
add_pic(s, "fig04_05_pruning_roc.png", Inches(7.9), Inches(1.35), width=Inches(5.1))
add_pic(s, "fig11_data_compare.png", Inches(7.9), Inches(4.0), width=Inches(5.1))

# ============ Slide 8: Fix 1 (Edges) ============
s = new_slide(8)
add_title_bar(s, "Fix 1 (Edges): Top-k Per Track Instead of a Global Threshold")
add_bullets(s, [
    "Per track, keep its top-k most confident edges (k\u224810\u201320) instead of a fixed 0.2 cut",
    "Hard compute cap: \u2264 N\u00d7k (~1.5k pairs) vs ~11k fully connected",
    "Inside top-k, keep low-confidence true edges \u2014 LCA classifier decides their class",
    "Why not \u201ckeep everything + soft weights\u201d: tree reconstruction is a discrete clustering \u2014 weights have no entry point; all background edges would merge the event into one cluster",
    "Pure decode change \u2014 testable on existing models immediately",
], Inches(0.8), Inches(1.35), Inches(7.3), Inches(5.5), size=15)
add_pic(s, "fig12_eb2_schematic.png", Inches(8.3), Inches(1.35), width=Inches(4.7))

# ============ Slide 9: Fix 2 (Nodes) ============
s = new_slide(9)
add_title_bar(s, "Fix 2 (Nodes): Train With a Smooth Pruning Mask")
add_bullets(s, [
    "Nodes and edges have very different scales (~150 vs ~23k/event) \u2192 need different fixes",
    "Make node keep/drop a smooth mask during training (temperature-annealed)",
    "The model \u201cexperiences\u201d pruning and becomes robust to it (closes train/inference gap)",
    "Fully vectorized tensor ops \u2192 GPU-parallel, negligible cost",
    "Motivation: on CERN data, node pruning alone kills ~26% of chains",
], Inches(0.8), Inches(1.35), Inches(11.5), Inches(5.0), size=18)

# ============ Slide 10: Fix 3 (Selection MLP) ============
s = new_slide(10)
add_title_bar(s, "Fix 3: Choose the Best Reconstruction Among Candidates")
add_bullets(s, [
    "The greedy reconstruction has no backtracking; softer candidates can yield several plausible trees per event",
    "Add a small MLP that scores each candidate tree (size, edge confidence, isolation from background) and selects the most plausible",
    "Needs one extra supervised training step (labels = truth-matching tree)",
], Inches(0.8), Inches(1.35), Inches(7.3), Inches(4.0), size=16)
add_pic(s, "fig12_eb2_schematic.png", Inches(8.3), Inches(1.35), width=Inches(4.7))

# ============ Slide 11: Expected Impact ============
s = new_slide(11)
add_title_bar(s, "Expected Impact")
add_bullets(s, [
    "Current (inactive-weights model, thr 0.2): 27.7% per-chain",
    "After retraining with corrected weights: ~42\u201350%",
    "Retrain + the three fixes: ~55\u201365% per-chain (\u224843\u201351% per-event)",
    "Caveats: estimates anchored on paper class accuracies; CERN lacks class-2/3; hyper-parameters not tuned",
], Inches(0.8), Inches(1.35), Inches(6.9), Inches(5.0), size=16)
add_pic(s, "fig08_expected_gain.png", Inches(7.9), Inches(1.35), width=Inches(5.1))

# ============ Slide 12: Next Steps ============
s = new_slide(12)
add_title_bar(s, "Next Steps")
add_bullets(s, [
    "Finish & evaluate the CERN retrain (v31); re-evaluate the public retrain (v30) at thr 0.2",
    "Implement the edge top-k change (decode-only) and test on existing models",
    "Put the node smooth-mask training + selection MLP into the next retrain",
    "Clarify with the group the track-sorting issue in the current CERN production data",
], Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0), size=19)

prs.save(OUT)
print("saved ->", OUT)
